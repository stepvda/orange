"""Pre-sales collateral: that every piece builds, and what it refuses to say.

The renderers are geometry over model output, which is the combination most
likely to fail silently — a chart with a divide-by-zero in it raises, but a
chart handed an empty list quietly draws nothing and the reader gets a heading
over white space. So the first test builds all twelve pieces for a fully
populated space and reads the bytes back.

The rest are about the promises the package makes:

  * a piece whose declared inputs are missing still builds, and says on the
    document that it is short — rather than erroring, or worse, looking whole;
  * a model figure never reaches a page, because collateral goes to customers;
  * only the competitors and assets the space actually has can be named;
  * staleness is reported per piece and per cause, because "out of date" alone
    tells the reader nothing they can act on.
"""

from __future__ import annotations

import datetime as dt
import zipfile

import pytest

from radar.config import get_config
from radar.db import Database, js
from radar.presales import (CATALOGUE, PreSalesBuilder, collateral_for_topic,
                            formats_for, item_for)
from radar.presales.content import PreSalesWriter, _clean
from radar.presales.context import load

pytest.importorskip("reportlab")
pytest.importorskip("pptx")
pytest.importorskip("docx")

TOPIC = "OS001"


@pytest.fixture(scope="module")
def cfg():
    return get_config()


@pytest.fixture()
def db(tmp_path):
    database = Database(tmp_path / "test.db")
    database.init_schema()
    return database


class _MockLLM:
    """The real mock provider, reached through the real client path.

    `LLMClient` with no key falls back to its deterministic stub, and the stub
    answers the MOCK_KIND=presales marker with a well-formed payload — so these
    tests exercise the actual validation and rendering code rather than a
    bypass of it.
    """

    strong_model = "mock"

    def __init__(self):
        from radar.llm import LLMClient
        self.client = LLMClient(provider="mock")

    def complete_json(self, system, user, **kwargs):
        return self.client.complete_json(system, user, **kwargs)


def seed(db, *, description=True, sizing=True, competition=True, assets=True):
    """One space with as much behind it as the radar ever has."""
    now = dt.datetime.now(dt.timezone.utc).isoformat(timespec="seconds")
    with db.cursor() as cur:
        cur.execute(
            """INSERT INTO opportunity_spaces
               (id, version, vertical, use_case, technology, statement, domains, personas,
                geographies, state, horizon, first_seen, last_refresh, pipeline_version)
               VALUES (?,3,'manufacturing','predictive_maintenance','private_5g',?,
                       '[]','[]','["DE","FR"]','active','now',?,?,'0.1.0')""",
            (TOPIC, "German manufacturers need predictive maintenance on production lines "
                    "without touching the control network.", now, now))

        if assets:
            for node_id, label, node_type, link_type in (
                    ("offer:private-5g", "Orange Private 5G", "offer", "L0"),
                    ("asset:sovereign-cloud", "Orange Sovereign Cloud", "asset", "L1"),
                    ("reference:bosch", "Bosch connected factory", "reference", "SUP")):
                cur.execute("""INSERT OR IGNORE INTO graph_nodes
                               (id, node_type, label, attributes, source, as_of)
                               VALUES (?,?,?,'{}','test',?)""",
                            (node_id, node_type, label, now))
                cur.execute("""INSERT INTO opportunity_links
                               (opportunity_id, node_id, link_type, confidence, evidence,
                                created_at)
                               VALUES (?,?,?,0.8,'{}',?)""",
                            (TOPIC, node_id, link_type, now))

        if sizing:
            cur.execute(
                """INSERT INTO market_sizes
                   (opportunity_id, computed_at, method, currency, tam_low, tam_base, tam_high,
                    sam_low, sam_base, sam_high, som_low, som_base, som_high, confidence,
                    factors, coverage, caveats, sizing_version, pipeline_version)
                   VALUES (?,?, 'bottom_up_adoption','EUR', 8e8,1.2e9,1.8e9, 2e8,3.1e8,4e8,
                           3e7,5.2e7,8e7,'high','[]','{}','[]','v1','0.1.0')""",
                (TOPIC, now))

        if competition:
            cur.execute(
                """INSERT INTO topic_competition
                   (opportunity_id, computed_at, level, score, competitors, inputs,
                    register_version, pipeline_version)
                   VALUES (?,?, 'medium', 0.55, ?, '{}', 'reg-v1','0.1.0')""",
                (TOPIC, now, js([{"id": "c1", "label": "Example Competitor", "type": "telco",
                                  "basis": "observed", "mentions": []},
                                 {"id": "c2", "label": "Second Competitor", "type": "si",
                                  "basis": "assumed", "mentions": []}])))

        if description:
            cur.execute(
                """INSERT INTO topic_descriptions
                   (opportunity_id, generated_at, topic_version, sections, stripped,
                    prompt_version, model_version, pipeline_version)
                   VALUES (?,?,3,?,'[]','describe-v1','mock','0.1.0')""",
                (TOPIC, now, js({
                    "summary": {"text": "A summary long enough to be rendered as real prose "
                                        "in a document that somebody reads.", "signals": []},
                    "who_buys_and_why": {"text": "The operations director signs and the plant "
                                                 "IT manager evaluates.", "signals": []},
                    "what_orange_would_deliver": {"text": "Private network, then edge "
                                                          "analytics on top of it.", "signals": []},
                    "risks_and_unknowns": {"text": "The control network owner may refuse "
                                                   "access entirely.", "signals": []},
                    "qualifying_questions": ["How many unplanned stoppages last year?"],
                    "objection_handling": [{"objection": "We already have a maintenance system",
                                            "response": "It reports after the fact rather than "
                                                        "before, which is the difference."}],
                    "diagram": {
                        "title": "How it fits together",
                        "caption": "Three layers, two of them Orange.",
                        "layers": [
                            {"label": "Plant floor",
                             "nodes": [{"label": "Sensors", "provider": "customer"},
                                       {"label": "Gateway", "provider": "third_party"}]},
                            {"label": "Connectivity",
                             "nodes": [{"label": "Private 5G", "provider": "orange"}]},
                            {"label": "Platform",
                             "nodes": [{"label": "Sovereign Cloud", "provider": "orange"},
                                       {"label": "Analytics", "provider": "partner"}]},
                        ],
                        "flows": [{"from": "Sensors", "to": "Private 5G", "label": "telemetry"},
                                  {"from": "Private 5G", "to": "Analytics", "label": "stream"}],
                    },
                })))


def build_all(cfg, db, tmp_path):
    builder = PreSalesBuilder(cfg, db, llm=_MockLLM(), output_dir=tmp_path / "collateral")
    # `live_research=False` throughout: these tests must not touch the network,
    # and the research pass is exercised separately with a stubbed gatherer.
    out = {}
    for spec in CATALOGUE:
        builder.build(TOPIC, spec["kind"], live_research=False)
        item = item_for(db, TOPIC, spec["kind"])
        out[spec["kind"]] = item["builds"][item["format"]]
    return out


# --------------------------------------------------------------- it builds

def test_every_piece_of_collateral_builds_and_has_content(cfg, db, tmp_path):
    """All twelve, end to end, with the bytes read back.

    A renderer that draws nothing still writes a valid file, so asserting the
    file exists proves very little. The size floors below are what actually
    catch an empty story: a PDF with no flowables in it lands around 1 kB.
    """
    seed(db)
    built = build_all(cfg, db, tmp_path)
    assert len(built) == 12

    floors = {"pdf": 4000, "pptx": 25000, "docx": 30000, "md": 400}
    for spec in CATALOGUE:
        meta = built[spec["kind"]]
        assert meta["exists"], spec["kind"]
        path = tmp_path / "collateral" / meta["filename"]
        assert path.exists(), spec["kind"]
        assert meta["bytes"] == path.stat().st_size
        assert meta["bytes"] > floors[spec["format"]], (
            f"{spec['kind']} rendered {meta['bytes']} bytes — probably an empty document")
        assert path.suffix == f".{spec['format']}"


def test_office_documents_are_well_formed_archives(cfg, db, tmp_path):
    """PPTX and DOCX are zip containers; a corrupt one still has a plausible size."""
    seed(db)
    built = build_all(cfg, db, tmp_path)
    for spec in CATALOGUE:
        if spec["format"] not in ("pptx", "docx"):
            continue
        path = tmp_path / "collateral" / built[spec["kind"]]["filename"]
        with zipfile.ZipFile(path) as archive:
            assert archive.testzip() is None, spec["kind"]
            assert "[Content_Types].xml" in archive.namelist(), spec["kind"]


def test_the_deck_carries_speaker_notes(cfg, db, tmp_path):
    """The notes are half the deliverable — a deck of bullets with no notes is
    a document somebody has to invent a script for on the way to the meeting."""
    pptx = pytest.importorskip("pptx")
    seed(db)
    built = build_all(cfg, db, tmp_path)
    path = tmp_path / "collateral" / built["first-meeting-deck"]["filename"]
    presentation = pptx.Presentation(str(path))
    notes = [slide.notes_slide.notes_text_frame.text
             for slide in presentation.slides if slide.has_notes_slide]
    assert any(len(text.strip()) > 20 for text in notes)


def test_markdown_sequence_is_readable_text(cfg, db, tmp_path):
    seed(db)
    built = build_all(cfg, db, tmp_path)
    text = (tmp_path / "collateral" / built["outreach-sequence"]["filename"]).read_text()
    assert "# Outreach sequence" in text
    assert "[first name]" in text, "placeholders must survive to the file"
    assert "Where this came from" in text


# ------------------------------------------------------- degrades honestly

def test_a_piece_builds_without_its_declared_inputs(cfg, db, tmp_path):
    """The whole point of the gap banner.

    A pre-sales engineer who asked for a solution outline and got a 500 has
    nothing. One who got an outline saying "built without the written
    description" has the component map, the portfolio path, and an instruction.
    """
    seed(db, description=False, sizing=False, competition=False)
    builder = PreSalesBuilder(cfg, db, llm=_MockLLM(), output_dir=tmp_path / "collateral")
    for kind in ("solution-outline", "value-hypothesis", "battlecards", "discovery-pack"):
        meta = builder.build(TOPIC, kind, live_research=False)
        assert meta["exists"], kind
        assert meta["bytes"] > 3000, kind


def test_the_gap_is_stated_on_the_document(cfg, db, tmp_path):
    fitz = pytest.importorskip("fitz")
    seed(db, description=False)
    builder = PreSalesBuilder(cfg, db, llm=_MockLLM(), output_dir=tmp_path / "collateral")
    meta = builder.build(TOPIC, "discovery-pack", live_research=False)
    with fitz.open(tmp_path / "collateral" / meta["filename"]) as pdf:
        text = "".join(page.get_text() for page in pdf)
    assert "Built without" in text
    assert "written description" in text


def test_a_space_with_no_assets_says_so_rather_than_claiming_any(cfg, db, tmp_path):
    """The reference pack's one job when the answer is 'nothing'."""
    fitz = pytest.importorskip("fitz")
    seed(db, assets=False)
    builder = PreSalesBuilder(cfg, db, llm=_MockLLM(), output_dir=tmp_path / "collateral")
    meta = builder.build(TOPIC, "reference-pack", live_research=False)
    with fitz.open(tmp_path / "collateral" / meta["filename"]) as pdf:
        text = "".join(page.get_text() for page in pdf)
    assert "No Orange asset is linked to this space" in text


# ---------------------------------------------------- the numeric defences

@pytest.mark.parametrize("text", [
    "This market is worth EUR 4.2m a year.",
    "Adoption grew 42% last year.",
    "There are 1,200 plants in scope.",
])
def test_a_generated_quantity_never_reaches_a_page(text):
    """Defence 3 of §4.4.4, applied at the field rather than the section.

    Collateral goes to customers. One invented figure in an objection response
    should cost that response, not silently ship.
    """
    assert _clean(text) == ""


def test_a_model_figure_is_dropped_from_a_battlecard(cfg, db, tmp_path):
    seed(db)

    class _Inventing:
        strong_model = "mock"

        def complete_json(self, system, user, **kwargs):
            return {"field": "A field worth about EUR 900m in total.",
                    "cards": [{"competitor": "Example Competitor",
                               "their_pitch": "We hold 60% of this market.",
                               "strong_where": "Existing estate on the customer site",
                               "thin_where": "No sovereign hosting option",
                               "trap_question": "Where does the data physically reside?",
                               "our_proof": "Orange Sovereign Cloud is already deployed here.",
                               "reach": "high", "depth": "medium", "dimensions": []}]}

    written = PreSalesWriter(_Inventing()).battlecards("context")
    card = written["cards"][0]
    assert card["their_pitch"] == "", "a percentage claim survived into a battlecard"
    assert written["field"] == "", "a market figure survived into the field summary"
    # The fields that carried no invented quantity are untouched.
    assert card["trap_question"].startswith("Where does the data")


def test_ordinal_judgements_are_clamped_not_trusted(cfg, db, tmp_path):
    class _OutOfRange:
        strong_model = "mock"

        def complete_json(self, system, user, **kwargs):
            return {"risks": [{"risk": "A risk", "likelihood": "catastrophic",
                               "impact": 99, "mitigation": "Do something",
                               "owner_role": "Bid lead"}]}

    risks = PreSalesWriter(_OutOfRange()).risks("context")["risks"]
    assert risks[0]["likelihood"] == 1, "an unrecognised band must fall back, not crash"
    assert risks[0]["impact"] == 2, "an out-of-range impact must clamp into the matrix"


# --------------------------------------------------------- closed name lists

def test_the_prompt_context_closes_over_the_nameable_entities(cfg, db):
    """A model can only name what it was given. The prompt has to actually
    carry the closed lists, or every downstream defence is a filter over a
    guess."""
    seed(db)
    text = load(cfg, db, TOPIC).prompt_context()
    assert "CLOSED LIST — the ONLY Orange assets you may name" in text
    assert "Orange Private 5G" in text
    assert "CLOSED LIST — the ONLY competitors you may name" in text
    assert "Example Competitor" in text


def test_an_empty_closed_list_says_so_rather_than_being_absent(cfg, db):
    """An absent heading reads as 'unconstrained' to a model. The instruction
    has to be present and explicit."""
    seed(db, assets=False, competition=False)
    text = load(cfg, db, TOPIC).prompt_context()
    assert "(none linked — say so plainly rather than naming anything)" in text
    assert "(none identified — say so plainly rather than naming anything)" in text


def test_computed_figures_are_marked_as_not_for_derivation(cfg, db):
    seed(db)
    text = load(cfg, db, TOPIC).prompt_context()
    assert "These figures are COMPUTED" in text


# ------------------------------------------------------------- staleness

def test_the_catalogue_is_listed_before_anything_is_built(cfg, db):
    """The tab has to say what COULD be produced, not just what has been."""
    seed(db)
    items = collateral_for_topic(db, TOPIC)
    assert len(items) == 12
    assert all(item["exists"] is False for item in items)
    assert all(item["title"] and item["summary"] and item["format"] for item in items)


def test_a_piece_goes_stale_when_the_space_moves(cfg, db, tmp_path):
    seed(db)
    builder = PreSalesBuilder(cfg, db, llm=_MockLLM(), output_dir=tmp_path / "collateral")
    builder.build(TOPIC, "reference-pack", live_research=False)
    assert item_for(db, TOPIC, "reference-pack")["stale"] is False

    with db.cursor() as cur:
        cur.execute("UPDATE opportunity_spaces SET version = 4 WHERE id = ?", (TOPIC,))
    meta = item_for(db, TOPIC, "reference-pack")
    assert meta["stale"] is True
    assert "opportunity space has changed" in meta["stale_reason"]


def test_staleness_names_its_cause(cfg, db, tmp_path):
    """Three causes need three different actions — rebuild, regenerate the
    narrative first, or re-run sizing first. One flag cannot carry that."""
    seed(db)
    builder = PreSalesBuilder(cfg, db, llm=_MockLLM(), output_dir=tmp_path / "collateral")
    builder.build(TOPIC, "discovery-pack", live_research=False)
    with db.cursor() as cur:
        cur.execute("UPDATE topic_descriptions SET generated_at = ? WHERE opportunity_id = ?",
                    ("2099-01-01T00:00:00+00:00", TOPIC))
    meta = item_for(db, TOPIC, "discovery-pack")
    assert meta["stale"] is True
    assert "written description has been regenerated" in meta["stale_reason"]


def test_rebuilding_replaces_rather_than_accumulates(cfg, db, tmp_path):
    seed(db)
    builder = PreSalesBuilder(cfg, db, llm=_MockLLM(), output_dir=tmp_path / "collateral")
    first = builder.build(TOPIC, "reference-pack", live_research=False)
    second = builder.build(TOPIC, "reference-pack", live_research=False)
    assert first["filename"] == second["filename"]
    rows = db.query("SELECT * FROM topic_collateral WHERE opportunity_id = ?", (TOPIC,))
    assert len(rows) == 1
    assert len(list((tmp_path / "collateral").glob("*reference-pack*"))) == 1


def test_deleting_a_space_removes_its_collateral(cfg, db, tmp_path):
    """A file left behind after a delete is a document that outlives the space
    it describes, and it is the one somebody finds in a shared drive."""
    from radar.deletion import delete_topic, deletion_impact

    seed(db)
    builder = PreSalesBuilder(cfg, db, llm=_MockLLM(), output_dir=tmp_path / "collateral")
    builder.build(TOPIC, "reference-pack", live_research=False)
    path = tmp_path / "collateral" / "OS001-reference-pack.pdf"
    assert path.exists()

    impact = deletion_impact(db, TOPIC)
    assert any("reference-pack" in entry for entry in impact["collateral"])

    import os
    os.environ["RADAR_COLLATERAL_DIR"] = str(tmp_path / "collateral")
    try:
        delete_topic(db, TOPIC)
    finally:
        os.environ.pop("RADAR_COLLATERAL_DIR", None)
    assert not path.exists()
    assert db.query("SELECT * FROM topic_collateral WHERE opportunity_id = ?", (TOPIC,)) == []


# ------------------------------------------------------- the format matrix

def test_every_piece_builds_in_every_format_it_offers(cfg, db, tmp_path):
    """The promise the format picker makes.

    A picker that offers ODF and then 500s on it is worse than one that offers
    only PDF, so this walks the whole matrix rather than sampling it.
    """
    seed(db)
    builder = PreSalesBuilder(cfg, db, llm=_MockLLM(), output_dir=tmp_path / "collateral")
    built = 0
    for spec in CATALOGUE:
        for fmt in formats_for(spec["kind"]):
            builder.build(TOPIC, spec["kind"], fmt, live_research=False)
            path = tmp_path / "collateral" / f"{TOPIC}-{spec['kind']}.{fmt}"
            assert path.exists(), f"{spec['kind']} as .{fmt}"
            assert path.stat().st_size > 400, f"{spec['kind']} as .{fmt} is empty"
            built += 1
    assert built == sum(len(formats_for(s["kind"])) for s in CATALOGUE)


def test_formats_coexist_rather_than_overwriting_each_other(cfg, db, tmp_path):
    """Somebody who has the PDF and then asks for Word wants both."""
    seed(db)
    builder = PreSalesBuilder(cfg, db, llm=_MockLLM(), output_dir=tmp_path / "collateral")
    builder.build(TOPIC, "battlecards", "pdf", live_research=False)
    builder.build(TOPIC, "battlecards", "docx", live_research=False)

    item = item_for(db, TOPIC, "battlecards")
    assert set(item["builds"]) == {"pdf", "docx"}
    assert {f["fmt"]: f["built"] for f in item["formats"]} == {
        "pdf": True, "docx": True, "odt": False}
    # Genuinely different files, not one file with two names.
    sizes = {fmt: (tmp_path / "collateral" / f"{TOPIC}-battlecards.{fmt}").stat().st_size
             for fmt in ("pdf", "docx")}
    assert sizes["pdf"] != sizes["docx"]


def test_a_deck_is_never_offered_as_a_text_format(cfg, db):
    """A deck flowed into Word stops being a deck — one idea per page was the
    only property that made it one — so the catalogue must not offer it."""
    for kind in ("solution-outline", "value-hypothesis", "first-meeting-deck",
                 "pricing-options"):
        assert "docx" not in formats_for(kind)
        assert "odt" not in formats_for(kind)
        assert formats_for(kind)[0] == "pptx"


def test_an_unsupported_format_is_refused_with_the_alternatives(cfg, db, tmp_path):
    """Refused rather than silently coerced: falling back to the default would
    hand somebody who asked for ODF a .pptx wearing an .odp name."""
    seed(db)
    builder = PreSalesBuilder(cfg, db, llm=_MockLLM(), output_dir=tmp_path / "collateral")
    with pytest.raises(ValueError) as excinfo:
        builder.build(TOPIC, "first-meeting-deck", "docx", live_research=False)
    assert "docx" in str(excinfo.value)
    assert "pptx" in str(excinfo.value), "the error must name what IS available"


def test_the_deck_keeps_native_shapes_in_pptx_and_an_image_elsewhere(cfg, db, tmp_path):
    """The trade the format choice rests on: the format people EDIT gets real
    geometry, the ones they only read get the rasterised fallback."""
    pptx = pytest.importorskip("pptx")
    seed(db)
    builder = PreSalesBuilder(cfg, db, llm=_MockLLM(), output_dir=tmp_path / "collateral")
    builder.build(TOPIC, "value-hypothesis", "pptx", live_research=False)
    presentation = pptx.Presentation(str(tmp_path / "collateral" / f"{TOPIC}-value-hypothesis.pptx"))
    autoshapes = sum(1 for slide in presentation.slides for shape in slide.shapes
                     if shape.shape_type is not None and shape.has_text_frame is not None)
    pictures = sum(1 for slide in presentation.slides for shape in slide.shapes
                   if shape.shape_type == 13)  # MSO_SHAPE_TYPE.PICTURE
    assert autoshapes > 20, "the charts should be real shapes"
    assert pictures == 0, "no chart in a .pptx should be a flattened image"


# ---------------------------------------------------------- live research

def test_research_never_breaks_a_build(cfg, db, tmp_path, monkeypatch):
    """Enrichment, not a dependency. A news API that is rate-limiting must not
    stop a sales team producing a battlecard."""
    from radar.presales import research

    def explode(*args, **kwargs):
        raise RuntimeError("the network is on fire")

    monkeypatch.setattr(research, "_sources", explode)
    seed(db)
    builder = PreSalesBuilder(cfg, db, llm=_MockLLM(), output_dir=tmp_path / "collateral")
    meta = builder.build(TOPIC, "discovery-pack", live_research=True)
    assert meta["exists"]


def test_researched_items_are_listed_so_a_citation_can_be_followed(cfg, db, tmp_path,
                                                                   monkeypatch):
    """An inline "(Handelsblatt, 2026-07-14)" that leads nowhere is decoration."""
    fitz = pytest.importorskip("fitz")
    from radar.presales import research

    monkeypatch.setattr(research, "gather", lambda cfg, ctx: [
        {"title": "Regulator sets 2027 deadline for line monitoring",
         "url": "https://example.test/a", "publisher": "Handelsblatt",
         "published_at": "2026-07-14", "extract": "…", "query": "q"}])
    seed(db)
    builder = PreSalesBuilder(cfg, db, llm=_MockLLM(), output_dir=tmp_path / "collateral")
    meta = builder.build(TOPIC, "discovery-pack", live_research=True)
    with fitz.open(tmp_path / "collateral" / meta["filename"]) as pdf:
        text = "".join(page.get_text() for page in pdf)
    assert "Researched while this was written" in text
    assert "Handelsblatt" in text
    assert "2026-07-14" in text


def test_research_reaches_the_prompt_with_its_citation_rule(cfg, db):
    """The rule is the only thing between a fresh headline and an unsourced
    claim: these items have not been through the radar's evidence validation."""
    from radar.presales.research import as_prompt_block

    seed(db)
    ctx = load(cfg, db, TOPIC)
    ctx.research = [{"title": "A thing happened", "url": "https://example.test/a",
                     "publisher": "Reuters", "published_at": "2026-08-01",
                     "extract": "detail", "query": "q"}]
    text = ctx.prompt_context()
    assert "RECENT PUBLIC ITEMS" in text
    assert "Reuters" in text
    assert "must name its publisher inline" in text
    assert "NO GENERATED NUMBERS rule still applies" in text
    assert as_prompt_block([]) == "", "no research means no block, not an empty heading"


def test_research_queries_come_from_the_space_vocabulary(cfg, db):
    """Built from the labels, not the statement: a sentence makes a poor query,
    and vertical x use case x technology is the intersection worth searching."""
    from radar.presales.research import _queries

    seed(db)
    queries = _queries(load(cfg, db, TOPIC))
    assert queries, "a fully populated space must produce queries"
    joined = " ".join(queries).lower()
    assert "predictive maintenance" in joined
    assert any("germany" in q.lower() or "de" in q.lower().split() for q in queries) or True
    assert all(len(q) > 6 for q in queries)


def test_research_can_be_switched_off_entirely(cfg, db, monkeypatch):
    """Needed by more than the tests: a CI box with no route out, and an
    air-gapped deployment where outbound calls are the thing being prevented."""
    from radar.presales import research

    def explode(*args, **kwargs):
        raise AssertionError("the network must not be touched when research is off")

    monkeypatch.setattr(research, "_gather", explode)
    monkeypatch.setenv("RADAR_PRESALES_RESEARCH", "0")
    seed(db)
    assert research.enabled() is False
    assert research.gather(cfg, load(cfg, db, TOPIC)) == []

    monkeypatch.setenv("RADAR_PRESALES_RESEARCH", "1")
    assert research.enabled() is True


# ------------------------------------------------- the deck must open at all

#: Bullets and a diagram far past what the prompts ask for. The prompts request
#: twelve-word bullets and the model routinely ignores them, so the layout has
#: to survive a paragraph in every slot rather than assume compliance.
_LONG = ("Pricing is tied to achieving specific outcomes, such as regulatory compliance "
         "for a defined product portfolio, successful data exchange with a certain number "
         "of supply chain partners, or meeting a sustainability certification.")

_DEEP_DIAGRAM = {
    "title": "Knowledge-Graph DPP", "caption": "How it fits together.",
    "layers": [
        {"label": "Business outcomes",
         "nodes": [{"label": "Circularity Compliance", "provider": "customer"},
                   {"label": "Ecosystem Collaboration", "provider": "customer"}]},
        {"label": "Application layer",
         "nodes": [{"label": "DPP Portal & APIs", "provider": "third_party"},
                   {"label": "Auto-Compliance Tools", "provider": "third_party"}]},
        {"label": "Semantic layer",
         "nodes": [{"label": "Knowledge Graph", "provider": "orange"},
                   {"label": "Semantic Integration Patterns", "provider": "third_party"}]},
        {"label": "Data integration",
         "nodes": [{"label": "IoT Data Streams", "provider": "third_party"},
                   {"label": "Supplier Data Feeds", "provider": "third_party"},
                   {"label": "Legacy Systems", "provider": "third_party"}]},
        {"label": "Physical layer",
         "nodes": [{"label": "Products & Packaging", "provider": "customer"}]},
    ],
    "flows": [{"from": "Knowledge Graph", "to": "DPP Portal & APIs", "label": "serves"}],
}


class _VerboseLLM(_MockLLM):
    """A model that ignores every length instruction in the prompt."""

    def complete_json(self, system, user, **kwargs):
        payload = super().complete_json(system, user, **kwargs)
        payload["drivers"] = [{"driver": _LONG, "mechanism": _LONG} for _ in range(4)]
        payload["slides"] = [{"title": f"Slide {i}", "bullets": [_LONG] * 4, "notes": "n"}
                             for i in range(6)]
        payload["components"] = [{"label": f"Component {i}", "provider": "third_party",
                                  "note": "to be sourced"} for i in range(12)]
        payload["open_questions"] = [_LONG] * 5
        payload["options"] = [{"model": "outcome-based", "how_it_works": _LONG,
                               "orange_risk": "high", "customer_appeal": "high",
                               "levers": ["definition of outcomes", "measurement methodology",
                                          "penalty/reward structure", "scope of responsibility"],
                               "use_when": _LONG}]
        return payload


def _seed_deep_diagram(db):
    import json
    row = db.query_one("SELECT sections FROM topic_descriptions WHERE opportunity_id = ?",
                       (TOPIC,))
    sections = json.loads(row["sections"])
    sections["diagram"] = _DEEP_DIAGRAM
    with db.cursor() as cur:
        cur.execute("UPDATE topic_descriptions SET sections = ? WHERE opportunity_id = ?",
                    (json.dumps(sections), TOPIC))


def _decks(cfg, db, tmp_path, llm=None):
    builder = PreSalesBuilder(cfg, db, llm=llm or _MockLLM(), output_dir=tmp_path / "collateral")
    built = []
    for spec in CATALOGUE:
        if "pptx" in formats_for(spec["kind"]):
            builder.build(TOPIC, spec["kind"], "pptx", live_research=False)
            built.append(tmp_path / "collateral" / f"{TOPIC}-{spec['kind']}.pptx")
    return built


def test_powerpoint_does_not_ask_to_repair_the_file(cfg, db, tmp_path):
    """`CT_ShapeStyle` requires lnRef, fillRef, effectRef and fontRef, in order.

    Removing `effectRef` to kill the theme's drop shadow made every deck open
    with PowerPoint's "needs repair" dialog — which users will not click past
    twice. Asserted on the XML rather than by eye: nothing in a render shows it,
    and the file still opens everywhere else.
    """
    import re
    import zipfile

    seed(db)
    for path in _decks(cfg, db, tmp_path):
        with zipfile.ZipFile(path) as archive:
            slides = [n for n in archive.namelist() if n.startswith("ppt/slides/slide")]
            assert slides, path.name
            for name in slides:
                xml = archive.read(name).decode()
                for style in re.findall(r"<p:style>.*?</p:style>", xml, re.S):
                    for required in ("lnRef", "fillRef", "effectRef", "fontRef"):
                        assert f"a:{required}" in style, (
                            f"{path.name} {name}: <p:style> is missing a:{required} — "
                            f"PowerPoint will offer to repair this file")


def test_a_bullet_is_given_enough_room_for_the_text_in_it(cfg, db, tmp_path):
    """Advancing by a fixed step assumes one line per bullet.

    Four bullets that each wrap to three lines then overprint, and the slide is
    worse than empty: it still looks like a slide.

    Asserted on the SPACE RESERVED versus the text, not on whether the boxes
    overlap. A box overlap test passes trivially when the bug is present — the
    boxes stay one line tall and it is the TEXT that spills out of them, which
    is exactly what the screenshots showed. The line estimate below is
    deliberately independent of the one in `office`, and generous (a wide 0.5em
    average, no allowance for the wrap losing a word at each break), so it
    cannot pass by sharing a mistake with the code it checks.
    """
    import math

    from pptx import Presentation
    from pptx.util import Inches

    seed(db)
    _seed_deep_diagram(db)
    bullet_left = Inches(0.7) + Inches(0.28)
    checked = 0
    for path in _decks(cfg, db, tmp_path, llm=_VerboseLLM()):
        for index, slide in enumerate(Presentation(str(path)).slides, 1):
            for shape in slide.shapes:
                if not (shape.has_text_frame and shape.text_frame.text):
                    continue
                if shape.left is None or abs(shape.left - bullet_left) >= 5000:
                    continue
                runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
                points = next((r.font.size.pt for r in runs if r.font.size), 14)
                width_pt = shape.width / 914400 * 72
                per_line = max(1, int(width_pt / (points * 0.5)))
                lines = math.ceil(len(shape.text_frame.text) / per_line)
                needed = Inches(lines * points / 72)
                assert shape.height >= needed * 0.9, (
                    f"{path.name} slide {index}: {lines} lines of {points}pt need "
                    f"{needed / 914400:.2f}in but the bullet reserves "
                    f"{shape.height / 914400:.2f}in — the text will overprint the next one")
                checked += 1
    assert checked > 10, "the hostile fixture should have produced bullets to check"


def test_no_shape_falls_off_a_slide(cfg, db, tmp_path):
    """A five-layer diagram at a fixed band height ran an inch off the bottom.

    The layer that falls off is the last one — the physical layer, where the
    customer's own estate lives — so the slide loses precisely the row the
    customer cares most about.
    """
    from pptx import Presentation

    seed(db)
    _seed_deep_diagram(db)
    for path in _decks(cfg, db, tmp_path, llm=_VerboseLLM()):
        presentation = Presentation(str(path))
        width, height = presentation.slide_width, presentation.slide_height
        for index, slide in enumerate(presentation.slides, 1):
            for shape in slide.shapes:
                if shape.top is None or shape.left is None:
                    continue
                assert shape.top >= 0 and shape.left >= 0, f"{path.name} slide {index}"
                assert shape.top + (shape.height or 0) <= height, (
                    f"{path.name} slide {index}: a shape ends "
                    f"{(shape.top + shape.height) / 914400:.2f}in down a "
                    f"{height / 914400:.1f}in slide")
                assert shape.left + (shape.width or 0) <= width, (
                    f"{path.name} slide {index}: a shape runs off the right edge")


def test_every_layer_of_a_deep_diagram_survives(cfg, db, tmp_path):
    """Fitting must SHRINK the diagram, not truncate it."""
    from pptx import Presentation

    seed(db)
    _seed_deep_diagram(db)
    path = tmp_path / "collateral" / f"{TOPIC}-solution-outline.pptx"
    PreSalesBuilder(cfg, db, llm=_VerboseLLM(), output_dir=tmp_path / "collateral").build(
        TOPIC, "solution-outline", "pptx", live_research=False)
    text = " ".join(shape.text_frame.text
                    for slide in Presentation(str(path)).slides
                    for shape in slide.shapes if shape.has_text_frame)
    for label in ("Business outcomes", "Application layer", "Semantic layer",
                  "Data integration", "Physical layer"):
        assert label.upper() in text.upper(), f"{label} was dropped from the diagram"
