"""PowerPoint: the house style, the native chart shapes, and the emitter.

WHY THE CHARTS ARE DRAWN AS NATIVE SHAPES rather than pasted in as images. A
picture in a deck as a PNG is a dead end: the architect who needs to move one
box, rename a component or delete the row that does not apply cannot, so they
redraw the whole slide — and what they redraw is no longer traceable to the
space it came from. Every diagram here is real PowerPoint geometry, rectangles
and connectors and freeforms at exact EMU coordinates, so it arrives editable,
scales without blurring and keeps its text searchable.

That costs a second implementation of geometry `charts` already has for PDF, and
it is worth paying exactly once — for the format people EDIT. The formats people
only read (Word, ODF) take the rasterised fallback in `emitters` instead, and no
third implementation is written.

Pure Python: no LibreOffice, no headless browser, nothing that would close the
sovereign deployment option (NFR-05).
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Any, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from ..sizing import format_eur
from .context import TopicContext

# ---------------------------------------------------------------------------
# Palette, restated in PowerPoint's colour type.
#
# The same validated values as `charts`: the ordinal Orange ramp for ordered
# stages, the all-pairs-safe categorical set for series that are merely
# different, warm-against-cool for signed quantities. A deck and the PDF beside
# it must not disagree about what orange means.
# ---------------------------------------------------------------------------

ORANGE = RGBColor.from_string("F16E00")
ORANGE_DARK = RGBColor.from_string("C24E00")
INK = RGBColor.from_string("141414")
MUTED = RGBColor.from_string("666666")
RULE = RGBColor.from_string("DDDDDD")
SURFACE = RGBColor.from_string("F6F6F6")
WHITE = RGBColor.from_string("FFFFFF")

ORDINAL = [RGBColor.from_string(v) for v in ("F08A2E", "DE6A05", "B85200", "8A3800")]
CATEGORICAL = [RGBColor.from_string(v) for v in ("D9600A", "2A78D6", "1BAF7A", "4A3AA7")]
GAIN, COST, NEUTRAL = CATEGORICAL[0], CATEGORICAL[1], RGBColor.from_string("9A9A9A")

PROVIDER_FILL = {
    "orange": (ORANGE, WHITE, ORANGE_DARK),
    "partner": (RGBColor.from_string("F9D3AE"), INK, ORANGE),
    "customer": (WHITE, INK, RGBColor.from_string("9A9A9A")),
    "third_party": (RGBColor.from_string("EFEFEF"), INK, RGBColor.from_string("C4C4C4")),
}
PROVIDER_LABEL = {"orange": "Orange asset", "partner": "Partner",
                  "customer": "Customer-owned", "third_party": "To be sourced"}

#: 16:9. Every geometry constant below is in inches against this frame.
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.7)
BODY_W = SLIDE_W - 2 * MARGIN


def _clamp(value: Any, low: float = 0.0, high: float = 1.0) -> float:
    try:
        return max(low, min(high, float(value)))
    except (TypeError, ValueError):
        return low


def _clip(text: Any, length: int) -> str:
    text = " ".join(str(text or "").split())
    return text if len(text) <= length else text[: length - 1] + "…"


# ---------------------------------------------------------------------------
# Shape hygiene
#
# Two things PowerPoint does by default that this house style does not want,
# both fixed in one place so no individual chart has to remember.
# ---------------------------------------------------------------------------

def fit(top: Emu, wanted: Emu, floor: float = 1.6) -> int:
    """`wanted` height, reduced to whatever is left below `top`.

    Charts used to take a default height chosen for a slide with no bullets on
    it. Once bullets were measured rather than assumed, `top` became variable
    and those defaults started running off the bottom of the slide. Every chart
    now asks for a height and is told what it can have.
    """
    available = float(SLIDE_H - top - Inches(0.6))
    return int(max(min(float(wanted), available), Inches(floor)))


def _flatten(shape) -> Any:
    """Strip the theme's drop shadow and 3-D effect from an autoshape.

    `shadow.inherit = False` alone is not enough: `add_shape` also writes a
    `<p:style>` block with an `effectRef` pointing into the theme, and a
    renderer that honours it puts a soft shadow under every box. The data is the
    only thing on these charts allowed to have weight.

    The reference is POINTED AT NOTHING rather than deleted. `CT_ShapeStyle`
    requires all four of lnRef, fillRef, effectRef and fontRef, in that order —
    removing one produces a file PowerPoint offers to repair before it will open
    it, which is a far worse outcome than a drop shadow. `idx="0"` is the
    schema's own way of saying "no entry from the theme's effect style list",
    so the element stays and the effect goes.
    """
    shape.shadow.inherit = False
    element = shape._element
    style = element.find(qn("p:style"))
    if style is not None:
        effect_ref = style.find(qn("a:effectRef"))
        if effect_ref is not None:
            effect_ref.set("idx", "0")
    return shape


def _text_width(text: str, size_pt: float, bold: bool = False) -> int:
    """Rendered width of a string in EMU, estimated.

    python-pptx cannot measure text — there is no font engine behind it — so a
    label that will not fit inside a bar cannot be detected the way the PDF
    charts detect it with `canvas.stringWidth`. This approximates Arial's
    average advance width and is deliberately PESSIMISTIC: over-estimating
    pushes a borderline label outside the bar, which is always readable, while
    under-estimating leaves white text overflowing a short bar onto a white
    slide, which is not readable at all.
    """
    per_char = 0.58 if bold else 0.52
    return int(len(text) * per_char * size_pt * 12700)


# ---------------------------------------------------------------------------
# Deck primitives
# ---------------------------------------------------------------------------

class Deck:
    """A 16:9 presentation with one house style, applied in one place."""

    def __init__(self, ctx: TopicContext, subject: str):
        self.ctx = ctx
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = SLIDE_W, SLIDE_H
        self.prs.core_properties.title = f"{ctx.topic_id} — {subject}"
        self.prs.core_properties.author = "Orange Business Innovation Radar"
        self.prs.core_properties.subject = ctx.statement[:200]
        self.subject = subject

    # -- chrome ------------------------------------------------------------

    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _rule(self, slide) -> None:
        """The Orange rule across the top of every slide — the same six points
        of colour the PDFs carry, so a deck and a brief read as one system."""
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.07))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ORANGE
        bar.line.fill.background()
        _flatten(bar)

    def _footer(self, slide) -> None:
        box = slide.shapes.add_textbox(MARGIN, SLIDE_H - Inches(0.45), BODY_W, Inches(0.3))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = f"{self.ctx.topic_id} · {_clip(self.ctx.statement, 96)}"
        run.font.size = Pt(9)
        run.font.color.rgb = MUTED

    def _text(self, slide, left, top, width, height, text: str, *, size: float,
              bold: bool = False, colour: RGBColor = INK, align=PP_ALIGN.LEFT,
              wrap: bool = True) -> Any:
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.word_wrap = wrap
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = colour
        # Helvetica is the PDFs' face; Arial is its metric-compatible sibling
        # and is the one that will actually be present on a sales laptop.
        run.font.name = "Arial"
        return box

    # -- slide kinds -------------------------------------------------------

    def cover(self, title: str, strapline: str) -> None:
        slide = self._blank()
        self._rule(slide)
        block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.35), Inches(0.16),
                                       Inches(2.1))
        block.fill.solid()
        block.fill.fore_color.rgb = ORANGE
        block.line.fill.background()
        _flatten(block)

        self._text(slide, MARGIN, Inches(2.3), BODY_W, Inches(0.5),
                   self.ctx.topic_id + " · " + self.ctx.triple, size=13, colour=MUTED)
        self._text(slide, MARGIN, Inches(2.85), BODY_W, Inches(1.4), title, size=36, bold=True)
        self._text(slide, MARGIN, Inches(4.2), BODY_W, Inches(0.8), self.ctx.statement, size=15,
                   colour=INK)
        self._text(slide, MARGIN, Inches(5.1), BODY_W, Inches(0.6), strapline, size=11,
                   colour=MUTED)
        self._footer(slide)

    def section(self, title: str, note: str = "") -> None:
        slide = self._blank()
        self._rule(slide)
        self._text(slide, MARGIN, Inches(3.0), BODY_W, Inches(0.9), title, size=28, bold=True)
        if note:
            self._text(slide, MARGIN, Inches(3.9), BODY_W, Inches(0.8), note, size=12, colour=MUTED)
        self._footer(slide)

    def content(self, title: str, bullets: Sequence[str], notes: str = "",
                subtitle: str = "") -> Any:
        """A titled slide. Returns it so a caller can draw a chart underneath.

        Bullets are advanced by their MEASURED height, not by a constant. A
        fixed 0.5in step assumes every bullet is one line; four bullets that each
        wrap to three lines then print on top of each other, and the slide is
        unreadable — worse than unreadable, because it still looks like a slide.
        The model is told to keep bullets to twelve words, but a prompt is a
        request and the layout has to survive it being ignored.
        """
        slide = self._blank()
        self._rule(slide)
        self._text(slide, MARGIN, Inches(0.45), BODY_W, Inches(0.7), title, size=24, bold=True)
        top = Inches(1.25)
        if subtitle:
            self._text(slide, MARGIN, Inches(1.12), BODY_W, Inches(0.4), subtitle, size=11,
                       colour=MUTED)
            top = Inches(1.62)

        shown = list(bullets[:6])
        text_w = BODY_W - Inches(0.28)
        size, line_h, gap = self._bullet_metrics(shown, text_w, float(SLIDE_H - top - Inches(0.9)))
        cursor = top
        for bullet in shown:
            lines = self._line_count(bullet, text_w, size)
            height = int(line_h * lines)
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, MARGIN, cursor + int(line_h * 0.34),
                Inches(0.09), Inches(0.09))
            marker.fill.solid()
            marker.fill.fore_color.rgb = ORANGE
            marker.line.fill.background()
            _flatten(marker)
            self._text(slide, MARGIN + Inches(0.28), cursor, text_w, height, bullet, size=size)
            cursor += height + gap

        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        self._footer(slide)
        # Where the bullets actually ended, so a caller placing a chart under
        # them does not have to re-derive it from a step size that is no longer
        # constant. Recomputing it at the call site is how the two drift apart.
        self.content_bottom = cursor + Inches(0.12)
        return slide

    @staticmethod
    def _line_count(text: str, width: int, size_pt: float) -> int:
        """How many lines this bullet will wrap to, estimated and never zero."""
        return max(1, -(-_text_width(text, size_pt) // max(int(width), 1)))

    def _bullet_metrics(self, bullets: Sequence[str], width: int,
                        available: float) -> tuple[float, float, int]:
        """Font size, line height and inter-bullet gap that fit the space left.

        Steps the type down rather than letting the block overflow the slide.
        14pt is the house size and where this starts; below 10pt a bullet stops
        being a bullet, so the last resort is accepting a tighter fit rather
        than shrinking further — the notes carry the detail either way.
        """
        for size in (14, 13, 12, 11, 10):
            # A point is 1/72in; 1.32 is the leading. The first version of this
            # used 0.052in per POINT, which made a 14pt line 0.73in tall — three
            # times reality — so every bullet reserved five lines' worth of slide
            # for three lines of text and the block ran out of room for no reason.
            line_h = int(Inches(1 / 72) * size * 1.32)
            gap = int(line_h * 0.55)
            total = sum(self._line_count(b, width, size) * line_h + gap for b in bullets)
            if total <= available:
                return size, line_h, gap
        line_h = int(Inches(1 / 72) * 10 * 1.32)
        return 10, line_h, int(line_h * 0.4)

    def save(self, path: Path) -> None:
        self.prs.save(str(path))


def _box(slide, left, top, width, height, label: str, fill: RGBColor, text_colour: RGBColor,
         border: RGBColor, note: str = "", size: float = 11) -> Any:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1)
    _flatten(shape)
    # The default rounded-rectangle radius is a bubble at this size; 0.08 of the
    # short side is the same restrained corner the PDFs use.
    shape.adjustments[0] = 0.08
    frame = shape.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = label
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_colour
    run.font.name = "Arial"
    if note:
        second = frame.add_paragraph()
        second.alignment = PP_ALIGN.CENTER
        run = second.add_run()
        run.text = note
        run.font.size = Pt(8)
        run.font.color.rgb = text_colour
        run.font.name = "Arial"
    return shape


def _legend(deck: Deck, slide, left, top, entries: Sequence[tuple[str, RGBColor]]) -> None:
    """Always present for two or more series, never for one."""
    if len(entries) < 2:
        return
    cursor = left
    for label, colour in entries:
        swatch = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cursor, top + Inches(0.03),
                                        Inches(0.14), Inches(0.11))
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = colour
        swatch.line.fill.background()
        _flatten(swatch)
        deck._text(slide, cursor + Inches(0.2), top - Inches(0.03), Inches(1.7), Inches(0.25),
                   label, size=9, colour=MUTED)
        cursor += Inches(0.2) + Inches(0.13 * len(label)) + Inches(0.2)


# ---------------------------------------------------------------------------
# Native chart shapes
# ---------------------------------------------------------------------------

def funnel(deck: Deck, slide, ctx: TopicContext, top: Emu) -> None:
    """TAM / SAM / SOM as nested bars on one baseline.

    Nested, not additive: SOM sits inside SAM sits inside TAM, and three bars
    drawn to one scale from one left edge is the only arrangement that says so.
    The ordinal ramp carries the narrowing, so the order survives without a
    legend.
    """
    size = ctx.best_size
    if not size:
        deck._text(slide, MARGIN, top, BODY_W, Inches(0.4),
                   "This space has not been sized. Run the sizing stage and rebuild.",
                   size=12, colour=MUTED)
        return
    stages = [("TAM", size.get("tam", {}).get("base"), "total addressable"),
              ("SAM", size.get("sam", {}).get("base"), "serviceable"),
              ("SOM", size.get("som", {}).get("base"), "realistically obtainable")]
    stages = [(a, b, c) for a, b, c in stages if b]
    if not stages:
        return
    peak = max(float(value) for _, value, _ in stages)
    label_w, note_w = Inches(0.9), Inches(2.1)
    track = BODY_W - label_w - note_w
    height = Inches(0.95)
    pitch = Inches(1.4)

    for index, (label, value, note) in enumerate(stages):
        y = top + pitch * index
        deck._text(slide, MARGIN, y + Inches(0.22), label_w, Inches(0.4), label, size=13,
                   bold=True, colour=MUTED)
        width = max(int(track * (float(value) / peak)), Emu(60000))
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN + label_w, y,
                                     width, height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = ORDINAL[min(index, len(ORDINAL) - 1)]
        bar.line.fill.background()
        _flatten(bar)
        # The radius adjustment is a fraction of the SHORT side, so a constant
        # here would round a wide bar and a narrow one differently. Scaling it
        # to the bar's own width keeps every corner the same physical size.
        bar.adjustments[0] = min(0.14, float(Inches(0.11)) / max(float(width), 1.0))

        money = format_eur(value)
        # SOM can be two percent of TAM, and its bar is then far too short for
        # its own label. Measured (approximately — see `_text_width`) and moved
        # outside when it will not fit: white text overflowing a short bar onto
        # a white slide is invisible, which is worse than an unlabelled bar.
        if _text_width(money, 14, bold=True) + Inches(0.3) < width:
            frame = bar.text_frame
            frame.word_wrap = False
            frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            paragraph = frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            run = paragraph.add_run()
            run.text = "  " + money
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = WHITE
            run.font.name = "Arial"
        else:
            deck._text(slide, MARGIN + label_w + width + Inches(0.12), y + Inches(0.22),
                       Inches(1.6), Inches(0.4), money, size=14, bold=True, colour=INK)

        deck._text(slide, MARGIN + label_w + track + Inches(0.15), y + Inches(0.26),
                   note_w, Inches(0.35), note, size=10, colour=MUTED)

    deck._text(slide, MARGIN, top + pitch * len(stages) + Inches(0.15), BODY_W, Inches(0.4),
               f"{size.get('method_label')} · {size.get('confidence')} confidence · per year. "
               f"Computed from stored components, not estimated.",
               size=9.5, colour=MUTED)


def _hairline(slide, left, top, width) -> Any:
    """A gridline: solid, one step off the surface, and recessive.

    Drawn as a thin filled rectangle rather than a connector because a connector
    inherits the theme line style, which is neither thin nor grey.
    """
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(left), int(top), int(width), Emu(9525))
    rule.fill.solid()
    rule.fill.fore_color.rgb = RULE
    rule.line.fill.background()
    return _flatten(rule)


def waterfall(deck: Deck, slide, steps: Sequence[tuple[str, float, str]], top: Emu,
              height: Emu = Inches(4.4)) -> None:
    """How a number is built, one signed step at a time.

    Gains and costs take the warm and cool poles; subtotals sit on the baseline
    in neutral grey because a subtotal has no sign. Each bar starts where the
    last one ended, and a connector runs from where one landed to where the next
    begins — without those the bars read as three unrelated quantities and the
    form loses the only thing it is for.
    """
    if not steps:
        return
    height = fit(top, height)
    running, cursor = [], 0.0
    for label, delta, kind in steps:
        if kind == "total":
            running.append((0.0, cursor, NEUTRAL, label, cursor))
        else:
            nxt = cursor + delta
            running.append((min(cursor, nxt), max(cursor, nxt),
                            GAIN if delta >= 0 else COST, label, delta))
            cursor = nxt
    peak = max((high for _, high, _, _, _ in running), default=1.0) or 1.0
    slot = int(BODY_W / len(running))
    # Capped in absolute terms as well as relative: three steps across a 16:9
    # slide gives a 4-inch slot, and a bar that fills it stops reading as a bar.
    bar_w = min(Inches(1.15), slot - Inches(0.4))
    plot_h = height - Inches(0.95)

    for frac in (0.0, 0.5, 1.0):
        _hairline(slide, MARGIN, top + plot_h - int(plot_h * frac), BODY_W)

    previous_landing = None
    for index, (low, high, colour, label, value) in enumerate(running):
        x = MARGIN + index * slot + (slot - bar_w) // 2
        y0 = top + plot_h - int(plot_h * (high / peak))
        bar_h = max(int(plot_h * ((high - low) / peak)), Emu(40000))

        # The connector is drawn first so the bar sits over its end.
        if previous_landing is not None:
            _hairline(slide, previous_landing[0], previous_landing[1],
                      x - previous_landing[0])

        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y0, bar_w, bar_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = colour
        bar.line.fill.background()
        _flatten(bar)
        bar.adjustments[0] = min(0.12, float(Inches(0.08)) / max(float(bar_w), 1.0))
        previous_landing = (x + bar_w, y0 if colour is not COST else y0 + bar_h)

        deck._text(slide, x - Inches(0.35), y0 - Inches(0.32), bar_w + Inches(0.7), Inches(0.3),
                   format_eur(abs(value)), size=11, bold=True, align=PP_ALIGN.CENTER)
        deck._text(slide, MARGIN + index * slot, top + plot_h + Inches(0.14), slot, Inches(0.55),
                   label, size=10, colour=MUTED, align=PP_ALIGN.CENTER)

    _legend(deck, slide, MARGIN, top + height - Inches(0.15),
            [("value created", GAIN), ("cost to serve", COST), ("net", NEUTRAL)])


def payback(deck: Deck, slide, cumulative: Sequence[float], top: Emu,
            height: Emu = Inches(4.0)) -> None:
    """Cumulative position over time as a real freeform, not a picture.

    One series, so no legend — the title names it. Everything here exists to
    locate one moment: where the line crosses zero. That gets the only marker
    and the only direct label; the rest is a baseline and two end ticks.
    """
    values = list(cumulative)
    if len(values) < 2:
        return
    height = fit(top, height)
    low, high = min(min(values), 0.0), max(max(values), 0.0)
    span = (high - low) or 1.0
    plot_w, plot_h = BODY_W - Inches(1.2), height - Inches(0.7)
    origin_x = MARGIN + Inches(1.2)

    def point(index: int) -> tuple[int, int]:
        x = origin_x + int(plot_w * (index / (len(values) - 1)))
        y = top + plot_h - int(plot_h * ((values[index] - low) / span))
        return int(x), int(y)

    zero_y = top + plot_h - int(plot_h * ((0.0 - low) / span))
    baseline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, origin_x, zero_y,
                                      plot_w, Emu(9525))
    baseline.fill.solid()
    baseline.fill.fore_color.rgb = MUTED
    baseline.line.fill.background()
    _flatten(baseline)
    deck._text(slide, MARGIN, zero_y - Inches(0.12), Inches(1.05), Inches(0.25),
               "break even", size=9, colour=MUTED, align=PP_ALIGN.RIGHT)

    start = point(0)
    builder = slide.shapes.build_freeform(start[0], start[1])
    builder.add_line_segments([point(i) for i in range(1, len(values))], close=False)
    line = builder.convert_to_shape()
    line.fill.background()
    line.line.color.rgb = ORANGE
    line.line.width = Pt(2.25)
    _flatten(line)

    crossing = next((i for i in range(1, len(values)) if values[i - 1] < 0 <= values[i]), None)
    if crossing is not None:
        cx, cy = point(crossing)
        radius = Inches(0.09)
        # A ring in the slide colour rather than a stroke, so the marker stays
        # legible where it sits on the line it belongs to.
        for size_, colour in ((radius + Inches(0.04), WHITE), (radius, ORANGE_DARK)):
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - size_, cy - size_,
                                         size_ * 2, size_ * 2)
            dot.fill.solid()
            dot.fill.fore_color.rgb = colour
            dot.line.fill.background()
            _flatten(dot)
        deck._text(slide, cx + Inches(0.15), cy - Inches(0.3), Inches(2.4), Inches(0.3),
                   f"breaks even in period {crossing}", size=11, bold=True, colour=ORANGE_DARK)

    for index, align in ((0, PP_ALIGN.LEFT), (len(values) - 1, PP_ALIGN.RIGHT)):
        x, _ = point(index)
        deck._text(slide, x - Inches(0.6), top + plot_h + Inches(0.12), Inches(1.2), Inches(0.28),
                   f"period {index}", size=9, colour=MUTED, align=PP_ALIGN.CENTER)


def component_map(deck: Deck, slide, components: Sequence[dict[str, str]], top: Emu,
                  columns: int = 4) -> None:
    """What the engagement is made of, coloured by who owns each piece.

    The same four-way provider encoding as the brief's solution diagram, because
    "who owns this" is the same question on both. The grey boxes are the point:
    each one is something somebody has to source before this is sellable, and
    they are countable at a glance.
    """
    if not components:
        return
    gap = Inches(0.18)
    box_w = int((BODY_W - gap * (columns - 1)) / columns)
    rows_needed = (min(len(components), 12) + columns - 1) // columns
    available = float(SLIDE_H - top - Inches(0.85))
    box_h = int(max(min(Inches(1.25),
                        (available - gap * (rows_needed - 1)) / max(rows_needed, 1)),
                    Inches(0.6)))
    for index, component in enumerate(components[:12]):
        row, column = divmod(index, columns)
        fill, text_colour, border = PROVIDER_FILL.get(
            component.get("provider", "third_party"), PROVIDER_FILL["third_party"])
        _box(slide, MARGIN + column * (box_w + gap), top + row * (box_h + gap),
             box_w, box_h, str(component.get("label", "")), fill, text_colour, border,
             note=str(component.get("note", "")), size=11)
    rows = (min(len(components), 12) + columns - 1) // columns
    _legend(deck, slide, MARGIN, top + rows * (box_h + gap) + Inches(0.08),
            [(PROVIDER_LABEL[k], PROVIDER_FILL[k][0])
             for k in ("orange", "partner", "customer", "third_party")])


def layered_diagram(deck: Deck, slide, diagram: dict[str, Any], top: Emu) -> None:
    """The brief's layered solution diagram, rebuilt as movable shapes.

    Same structure the model produced and the same validation behind it; only
    the renderer differs. Flows are real connectors, so a box dragged in
    PowerPoint keeps its arrows.
    """
    layers = ((diagram or {}).get("layers") or [])[:5]
    if not layers:
        return
    label_w = Inches(1.5)
    usable = BODY_W - label_w
    placed: dict[str, Any] = {}

    # Sized to the space that is actually left, not to a constant. A five-layer
    # diagram at a fixed 1.15in band ran 1in off the bottom of the slide — and a
    # diagram that overflows is worse than a small one, because the layer that
    # falls off the page is the physical layer, which is where the customer's
    # own estate lives. The band shrinks to fit; below the floor it stops
    # shrinking and the legend is dropped instead.
    available = float(SLIDE_H - top - Inches(0.75))
    gap = Inches(0.18)
    layer_h = int(min(Inches(1.15),
                      (available - gap * (len(layers) - 1)) / max(len(layers), 1)))
    layer_h = max(layer_h, int(Inches(0.62)))
    node_h = int(min(Inches(0.78), layer_h - Inches(0.22)))

    for index, layer in enumerate(layers):
        y = top + index * (layer_h + gap)
        band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, y, BODY_W, layer_h)
        band.fill.solid()
        band.fill.fore_color.rgb = SURFACE
        band.line.color.rgb = RULE
        band.line.width = Pt(0.75)
        _flatten(band)
        band.adjustments[0] = 0.05
        deck._text(slide, MARGIN + Inches(0.12), y + (layer_h - Inches(0.3)) // 2,
                   label_w - Inches(0.2), Inches(0.4),
                   str(layer.get("label", "")).upper(), size=9, bold=True, colour=MUTED)

        nodes = (layer.get("nodes") or [])[:4]
        if not nodes:
            continue
        node_gap = Inches(0.14)
        node_w = int((usable - Inches(0.2) - node_gap * (len(nodes) - 1)) / len(nodes))
        for position, node in enumerate(nodes):
            fill, text_colour, border = PROVIDER_FILL.get(
                node.get("provider", "third_party"), PROVIDER_FILL["third_party"])
            shape = _box(slide, MARGIN + label_w + position * (node_w + node_gap),
                         y + (layer_h - node_h) // 2, node_w, node_h,
                         str(node.get("label", "")), fill, text_colour, border, size=10)
            placed[str(node.get("label", ""))] = shape

    for flow in ((diagram or {}).get("flows") or [])[:8]:
        source, target = placed.get(flow.get("from", "")), placed.get(flow.get("to", ""))
        if source is None or target is None or source is target:
            continue
        connector = slide.shapes.add_connector(
            2, int(source.left + source.width / 2), int(source.top),
            int(target.left + target.width / 2), int(target.top + target.height))
        connector.line.color.rgb = ORANGE
        connector.line.width = Pt(1.25)

    legend_y = top + len(layers) * (layer_h + gap) + Inches(0.04)
    if legend_y + Inches(0.25) < SLIDE_H - Inches(0.5):
        _legend(deck, slide, MARGIN, legend_y,
                [(PROVIDER_LABEL[k], PROVIDER_FILL[k][0])
                 for k in ("orange", "partner", "customer", "third_party")])


def option_columns(deck: Deck, slide, options: Sequence[dict[str, Any]], top: Emu,
                   measure: str, key: str, height: Emu = Inches(3.8)) -> None:
    """A handful of commercial models on one ordinal measure.

    Four slots is where the categorical palette stops, and it is also where a
    comparison stops being readable — a fifth option means this should be a
    table, not a chart.
    """
    if not options:
        return
    height = fit(top, height)
    slot = int(BODY_W / len(options))
    bar_w = min(Inches(1.5), slot - Inches(0.5))
    plot_h = height - Inches(0.6)
    words = {0: "low", 1: "medium", 2: "high"}
    for index, option in enumerate(options[:4]):
        band = int(option.get(key, 1))
        x = MARGIN + index * slot + (slot - bar_w) // 2
        bar_h = max(int(plot_h * ((band + 1) / 3)), Emu(50000))
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top + plot_h - bar_h,
                                     bar_w, bar_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = CATEGORICAL[index]
        bar.line.fill.background()
        _flatten(bar)
        bar.adjustments[0] = 0.1
        deck._text(slide, x - Inches(0.2), top + plot_h - bar_h - Inches(0.3),
                   bar_w + Inches(0.4), Inches(0.28), words[band], size=11, bold=True,
                   align=PP_ALIGN.CENTER)
        deck._text(slide, MARGIN + index * slot, top + plot_h + Inches(0.1), slot, Inches(0.5),
                   str(option.get("model", "")), size=11, align=PP_ALIGN.CENTER, colour=MUTED)
    deck._text(slide, MARGIN, top + height - Inches(0.15), BODY_W, Inches(0.3), measure,
               size=9, colour=MUTED)


def field_map(deck: Deck, slide, entries: Sequence[dict[str, Any]], top: Emu,
              height: Emu = Inches(4.4)) -> None:
    """Where each competitor stands, and where Orange stands.

    Not coloured by competitor: identity comes from the name printed beside the
    dot, which frees the one colour that matters here for Orange's own mark. It
    also sidesteps the series cap that a scatter imposes — any two dots can end
    up adjacent, so a per-competitor palette would have to separate every pair,
    and names do not have that limit.
    """
    if not entries:
        return
    height = fit(top, height)
    plot_w = BODY_W - Inches(1.0)
    plot_h = height - Inches(0.7)
    origin_x = MARGIN + Inches(0.75)
    plot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, origin_x, top, plot_w, plot_h)
    plot.fill.solid()
    plot.fill.fore_color.rgb = SURFACE
    plot.line.color.rgb = RULE
    plot.line.width = Pt(0.75)
    _flatten(plot)

    for entry in entries:
        orange = bool(entry.get("is_orange"))
        radius = Inches(0.11) if orange else Inches(0.085)
        cx = origin_x + int(plot_w * _clamp(entry.get("x", 0.5)))
        cy = top + plot_h - int(plot_h * _clamp(entry.get("y", 0.5)))
        for size_, colour in ((radius + Inches(0.035), WHITE),
                              (radius, ORANGE if orange else CATEGORICAL[1])):
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - size_, cy - size_,
                                         size_ * 2, size_ * 2)
            dot.fill.solid()
            dot.fill.fore_color.rgb = colour
            dot.line.fill.background()
            _flatten(dot)
        # Flip the label inboard near the right edge so it never leaves the plot.
        near_edge = cx > origin_x + plot_w - Inches(1.9)
        deck._text(slide,
                   cx - Inches(2.0) if near_edge else cx + Inches(0.16),
                   cy - Inches(0.13), Inches(1.85), Inches(0.28),
                   _clip(str(entry.get("label", "")), 26), size=10, bold=orange,
                   colour=ORANGE_DARK if orange else INK,
                   align=PP_ALIGN.RIGHT if near_edge else PP_ALIGN.LEFT)

    deck._text(slide, origin_x, top + plot_h + Inches(0.12), plot_w, Inches(0.3),
               "reach across this market →", size=9, colour=MUTED, align=PP_ALIGN.CENTER)
    deck._text(slide, MARGIN - Inches(0.1), top + plot_h / 2 - Inches(0.15), Inches(0.85),
               Inches(0.3), "depth ↑", size=9, colour=MUTED, align=PP_ALIGN.RIGHT)


def provenance_slide(deck: Deck, ctx: TopicContext, written: bool) -> None:
    """The last slide, and the one that makes the deck checkable six months on."""
    size = ctx.best_size
    rows = [
        ("Opportunity space", f"{ctx.topic_id} v{ctx.topic.get('version')}"),
        ("Weight set", str((ctx.topic.get("provenance") or {}).get("weight_set") or "—")),
        ("Sizing", f"{size.get('sizing_version')} ({size.get('method')})" if size else "not sized"),
        ("Competitor register", str((ctx.analysis or {}).get("register_version")
                                    or (ctx.competition or {}).get("register_version") or "—")),
    ]
    if written:
        rows.append(("Written sections",
                     "generated under the §4.4.4 defences — quantities stripped, not trusted"))
    slide = deck._blank()
    deck._rule(slide)
    deck._text(slide, MARGIN, Inches(0.45), BODY_W, Inches(0.6), "Where this came from",
               size=22, bold=True)
    deck._text(slide, MARGIN, Inches(1.1), BODY_W, Inches(0.4),
               "Every figure in this deck decomposes into stored components. Nothing on the "
               "money slides was written by a model.", size=11, colour=MUTED)
    for index, (key, value) in enumerate(rows):
        y = Inches(1.75) + Inches(0.46) * index
        deck._text(slide, MARGIN, y, Inches(2.6), Inches(0.4), key, size=11, bold=True,
                   colour=MUTED)
        deck._text(slide, MARGIN + Inches(2.7), y, BODY_W - Inches(2.7), Inches(0.4),
                   value, size=11)
    deck._footer(slide)


# ---------------------------------------------------------------------------
# The PowerPoint emitter
# ---------------------------------------------------------------------------

def deck_to_pptx(model: Any, ctx: TopicContext, path: Path) -> None:
    """A `blocks.Deck` as a .pptx, with charts as native shapes.

    The native path is the whole point of choosing PowerPoint: `Chart.pptx` is
    a builder that lays down real rectangles, connectors and freeforms, so the
    architect who needs to move one box can. A chart with no native builder
    falls back to the rasterised image, which is still better than omitting it —
    but every chart the decks actually use has one.
    """
    deck = Deck(ctx, model.subject)
    for slide_model in model.slides:
        if slide_model.kind == "cover":
            deck.cover(slide_model.title, slide_model.strapline)
            continue
        if slide_model.kind == "section":
            deck.section(slide_model.title, slide_model.subtitle)
            continue

        slide = deck.content(slide_model.title, slide_model.bullets,
                             notes=slide_model.notes, subtitle=slide_model.subtitle)
        top = deck.content_bottom

        for index, (key, value) in enumerate(slide_model.rows):
            y = top + Inches(0.46) * index
            deck._text(slide, MARGIN, y, Inches(2.6), Inches(0.4), key, size=11, bold=True,
                       colour=MUTED)
            deck._text(slide, MARGIN + Inches(2.7), y, BODY_W - Inches(2.7), Inches(0.4),
                       value, size=11)

        chart = slide_model.chart
        if chart is None:
            continue
        if chart.pptx is not None:
            chart.pptx(deck, slide, top)
        else:
            width_mm, height_mm = chart.size_mm()
            available = float(SLIDE_H - top - Inches(0.7))
            scale = min(float(BODY_W) / Inches(width_mm / 25.4),
                        available / max(float(Inches(height_mm / 25.4)), 1.0), 1.0)
            slide.shapes.add_picture(
                io.BytesIO(chart.png()), MARGIN, top,
                width=int(Inches(width_mm / 25.4) * scale),
                height=int(Inches(height_mm / 25.4) * scale))

    deck.save(path)
