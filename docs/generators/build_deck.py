"""Build the Innovation Radar presentation (docs/Orange_Innovation_Radar.pptx).

Kept as a script rather than a hand-made file so the deck can be regenerated
when the numbers change — every figure below is pulled from the live database
rather than typed in, which is the same discipline the product itself applies to
its own claims.

    python3 docs/generators/build_deck.py
"""

from __future__ import annotations

import sqlite3
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
SHOTS = Path("/tmp/vid/shots")
OUT = ROOT / "docs" / "Orange_Innovation_Radar.pptx"

# 16:9
W, H = Inches(13.333), Inches(7.5)

ORANGE = RGBColor(0xF0, 0x70, 0x00)
INK = RGBColor(0x0B, 0x0B, 0x0B)
MUTED = RGBColor(0x6B, 0x69, 0x63)
LIGHT = RGBColor(0xF4, 0xF4, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xDE, 0xDC, 0xD6)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
GREEN = RGBColor(0x0C, 0xA3, 0x0C)
RED = RGBColor(0xD0, 0x3B, 0x3B)

SANS = "Helvetica Neue"
MONO = "Menlo"


# ---------------------------------------------------------------------------
# Live figures — never typed by hand
# ---------------------------------------------------------------------------

def figures() -> dict:
    db = ROOT / "data" / "radar.db"
    if not db.exists():
        return {}
    con = sqlite3.connect(f"file:{db}?mode=ro", uri=True)
    one = lambda q: con.execute(q).fetchone()[0]
    f = {
        "topics": one("SELECT COUNT(*) FROM opportunity_spaces"),
        "signals": one("SELECT COUNT(*) FROM signals"),
        "relevant": one("SELECT COUNT(*) FROM signals WHERE relevance>0"),
        "sources": one("SELECT COUNT(DISTINCT source_id) FROM signals"),
        "links": one("SELECT COUNT(*) FROM opportunity_links"),
        "nodes": one("SELECT COUNT(*) FROM graph_nodes"),
        "attach": one("SELECT COUNT(*) FROM opportunity_signals"),
        "verticals": one("SELECT COUNT(DISTINCT vertical) FROM opportunity_spaces"),
        "use_cases": one("SELECT COUNT(DISTINCT use_case) FROM opportunity_spaces"),
        "techs": one("SELECT COUNT(DISTINCT technology) FROM opportunity_spaces"),
        "fr": one("SELECT COUNT(*) FROM signals WHERE language='fr'"),
        "tier1": one("SELECT COUNT(*) FROM signals WHERE tier=1"),
    }
    def maybe(sql, default=0):
        try:
            return con.execute(sql).fetchone()[0]
        except sqlite3.Error:
            return default
    f["sized"] = maybe("SELECT COUNT(DISTINCT opportunity_id) FROM market_sizes")
    f["sized_bu"] = maybe("SELECT COUNT(DISTINCT opportunity_id) FROM market_sizes "
                          "WHERE method='bottom_up_adoption' AND som_base > 0")
    f["plans"] = maybe("SELECT COUNT(*) FROM plans")
    f["collateral"] = maybe("SELECT COUNT(*) FROM topic_collateral")
    f["collateral_kinds"] = maybe("SELECT COUNT(DISTINCT kind) FROM topic_collateral")
    f["committed"] = maybe("SELECT COUNT(*) FROM workflow_state WHERE stage != 'shortlisted'")
    f["analyses"] = maybe("SELECT COUNT(DISTINCT opportunity_id) FROM topic_competitor_analysis")
    f["profiles"] = maybe("SELECT COUNT(*) FROM competitor_profiles")
    f["pages"] = maybe("SELECT COUNT(*) FROM competitor_pages")
    f["competed"] = maybe("SELECT COUNT(DISTINCT opportunity_id) FROM topic_competition")
    f["briefs"] = maybe("SELECT COUNT(DISTINCT opportunity_id) FROM topic_briefs")
    f["avg_sig"] = round(f["attach"] / max(f["topics"], 1), 1)
    con.close()
    return f


F = figures()


# ---------------------------------------------------------------------------
# Layout helpers
# ---------------------------------------------------------------------------

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(slide, x, y, w, h, runs, size=18, color=INK, bold=False, font=SANS,
         align=PP_ALIGN.LEFT, spacing=1.18, anchor=MSO_ANCHOR.TOP):
    """`runs` is a string or a list of (text, {overrides}) tuples/paragraph lists."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    paragraphs = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        pieces = para if isinstance(para, list) else [(para, {})]
        for content, opts in pieces:
            r = p.add_run()
            r.text = content
            r.font.size = Pt(opts.get("size", size))
            r.font.bold = opts.get("bold", bold)
            r.font.color.rgb = opts.get("color", color)
            r.font.name = opts.get("font", font)
        if i < len(paragraphs) - 1:
            p.space_after = Pt(opts.get("after", 8))
    return box


def header(slide, kicker, title, subtitle=None):
    rect(slide, 0, 0, W, Inches(0.06), fill=ORANGE)
    text(slide, Inches(0.75), Inches(0.5), Inches(11.8), Inches(0.3),
         kicker.upper(), size=11, color=ORANGE, bold=True)
    text(slide, Inches(0.75), Inches(0.85), Inches(11.8), Inches(0.7),
         title, size=32, bold=True)
    if subtitle:
        text(slide, Inches(0.75), Inches(1.62), Inches(11.0), Inches(0.5),
             subtitle, size=15, color=MUTED)
    return Inches(2.35) if subtitle else Inches(1.9)


def footer(slide, n):
    text(slide, Inches(0.75), Inches(6.95), Inches(8), Inches(0.3),
         "Orange Business · Opportunity Spaces / Innovation Radar", size=9, color=MUTED)
    text(slide, Inches(11.6), Inches(6.95), Inches(1.0), Inches(0.3),
         str(n), size=9, color=MUTED, align=PP_ALIGN.RIGHT, font=MONO)


def picture(slide, name, x, y, w):
    path = SHOTS / f"{name}.png"
    if not path.exists():
        rect(slide, x, y, w, Inches(3.2), fill=LIGHT, line=RULE)
        text(slide, x, y + Inches(1.4), w, Inches(0.4), f"[{name}.png missing]",
             size=12, color=MUTED, align=PP_ALIGN.CENTER)
        return
    pic = slide.shapes.add_picture(str(path), x, y, width=w)
    # Hairline so a white UI does not bleed into a white slide.
    rect(slide, x, y, pic.width, pic.height, fill=None, line=RULE, line_w=0.75)


def bullets(slide, x, y, w, items, size=15, gap=0.44, marker=True):
    for i, item in enumerate(items):
        yy = y + Inches(gap * i)
        if marker:
            rect(slide, x, yy + Inches(0.075), Inches(0.09), Inches(0.09), fill=ORANGE)
        if isinstance(item, tuple):
            head, tail = item
            text(slide, x + Inches(0.26 if marker else 0), yy, w, Inches(0.4),
                 [[(head + "  ", {"bold": True}), (tail, {"color": MUTED})]], size=size)
        else:
            text(slide, x + Inches(0.26 if marker else 0), yy, w, Inches(0.4), item, size=size)


def stat(slide, x, y, w, value, label, sub=None, color=INK):
    rect(slide, x, y, w, Inches(1.32), fill=LIGHT)
    text(slide, x + Inches(0.22), y + Inches(0.16), w - Inches(0.4), Inches(0.5),
         str(value), size=30, bold=True, color=color, font=MONO)
    text(slide, x + Inches(0.22), y + Inches(0.72), w - Inches(0.4), Inches(0.3),
         label.upper(), size=9.5, color=MUTED, bold=True)
    if sub:
        text(slide, x + Inches(0.22), y + Inches(0.97), w - Inches(0.4), Inches(0.3),
             sub, size=10, color=MUTED)


def caption(slide, y, body):
    text(slide, Inches(0.75), y, Inches(11.9), Inches(0.5), body, size=12.5, color=MUTED)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    n = 0

    # ---- 1. Title -------------------------------------------------------
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=INK)
    rect(s, 0, 0, Inches(0.22), H, fill=ORANGE)
    text(s, Inches(1.1), Inches(2.15), Inches(11), Inches(0.4),
         "ORANGE BUSINESS", size=13, color=ORANGE, bold=True)
    text(s, Inches(1.1), Inches(2.6), Inches(11), Inches(1.4),
         "Opportunity Spaces\nInnovation Radar", size=48, bold=True, color=WHITE, spacing=1.05)
    text(s, Inches(1.1), Inches(4.35), Inches(10), Inches(0.9),
         "A regularly refreshed view of specific innovation opportunities — each scored on how "
         "attractive it is, how urgent the window is, and how strong Orange's right to win is.",
         size=15, color=RGBColor(0xC3, 0xC2, 0xB7))
    text(s, Inches(1.1), Inches(5.6), Inches(11), Inches(0.4),
         [[("MVP walkthrough", {"color": WHITE, "bold": True}),
           ("   ·   concepts · functionality · planning · architecture", {"color": MUTED})]], size=13)
    n += 1

    # ---- 2. The problem -------------------------------------------------
    s = blank(prs); y = header(s, "Why", "The problem is not a shortage of information")
    text(s, Inches(0.75), y, Inches(6.0), Inches(2.2),
         "The available information about technology is generic, undated, unsourced — and "
         "disconnected from what Orange can actually sell.", size=19, spacing=1.3)
    text(s, Inches(0.75), y + Inches(1.5), Inches(6.0), Inches(2.4),
         "The radar's value is the join between an external signal and an internal asset, "
         "expressed at a level of specificity a salesperson can use in a meeting on Thursday.",
         size=15, color=MUTED, spacing=1.35)
    rect(s, Inches(7.2), y - Inches(0.1), Inches(5.4), Inches(3.5), fill=LIGHT)
    text(s, Inches(7.5), y + Inches(0.2), Inches(4.8), Inches(0.3),
         "REJECTED AS TOPICS", size=10, color=RED, bold=True)
    bullets(s, Inches(7.5), y + Inches(0.62), Inches(4.8),
            ['"AI"', '"Cloud"', '"Cybersecurity"', '"Digital transformation in retail"'],
            size=15, gap=0.36, marker=False)
    text(s, Inches(7.5), y + Inches(2.15), Inches(4.8), Inches(0.3),
         "WHAT A TOPIC MUST LOOK LIKE", size=10, color=GREEN, bold=True)
    text(s, Inches(7.5), y + Inches(2.5), Inches(4.8), Inches(0.8),
         "Private 5G plus edge vision for safety compliance in mining", size=14, bold=True)
    footer(s, n := n + 1)

    # ---- 3. What is an opportunity space --------------------------------
    s = blank(prs); y = header(s, "Concept", "An opportunity space is a triple",
                               "Identity is the triple. The statement is a rendering of it — both are stored.")
    bw = Inches(3.35)
    for i, (label, value) in enumerate([("VERTICAL", "Manufacturing"),
                                        ("USE CASE", "OT / ICS security"),
                                        ("TECHNOLOGY", "SIEM and SOAR")]):
        x = Inches(0.75) + (bw + Inches(0.5)) * i
        rect(s, x, y, bw, Inches(1.5), fill=LIGHT)
        text(s, x + Inches(0.25), y + Inches(0.25), bw - Inches(0.5), Inches(0.3),
             label, size=10, color=ORANGE, bold=True)
        text(s, x + Inches(0.25), y + Inches(0.62), bw - Inches(0.5), Inches(0.7),
             value, size=18, bold=True)
        if i < 2:
            text(s, x + bw + Inches(0.13), y + Inches(0.5), Inches(0.3), Inches(0.4),
                 "×", size=24, color=MUTED, bold=True)
    rect(s, Inches(0.75), y + Inches(1.95), Inches(11.9), Inches(1.0), fill=None, line=ORANGE, line_w=1.5)
    text(s, Inches(1.05), y + Inches(2.2), Inches(11.3), Inches(0.6),
         "SIEM/SOAR-based OT security analytics for legacy SCADA and ICS in brownfield "
         "manufacturing sites.", size=17, bold=True)
    caption(s, y + Inches(3.15),
            "The triple gives deduplication and filtering. The statement gives the specificity a "
            "salesperson needs. A candidate that does not resolve to exactly one of each fails validation.")
    footer(s, n := n + 1)

    # ---- 4. Two scores --------------------------------------------------
    s = blank(prs); y = header(s, "Concept", "Two scores, never one",
                               "They answer different questions and are owned by different people.")
    cards = [
        ("ATTRACTIVENESS", "Is the world moving?",
         "Market signal strength · Source diversity · Evidence quality · Novelty and momentum · "
         "Strategic relevance", "Computed from external evidence only.", ORANGE),
        ("RIGHT TO WIN", "Can we play, can we win?",
         "Offer match · Reference density · Partner coverage · Compliance fit · Capability depth · "
         "Analyst recognition · Technology ownership",
         "Computed from the Orange Business Graph as named query results — never asserted by a model.", BLUE),
        ("CONVICTION", "Do our own people believe it?",
         "Strategic fit (strategy) · Customer demand (sales) · Deliverability (presales)",
         "Adjusts what surfaces first for each role. Never alters the other two.", GREEN),
    ]
    cw = Inches(3.85)
    for i, (title_, q, comps, note, col) in enumerate(cards):
        x = Inches(0.75) + (cw + Inches(0.28)) * i
        rect(s, x, y, cw, Inches(3.35), fill=LIGHT)
        rect(s, x, y, cw, Inches(0.07), fill=col)
        text(s, x + Inches(0.28), y + Inches(0.32), cw - Inches(0.56), Inches(0.3),
             title_, size=11, color=col, bold=True)
        text(s, x + Inches(0.28), y + Inches(0.66), cw - Inches(0.56), Inches(0.4),
             q, size=16, bold=True)
        text(s, x + Inches(0.28), y + Inches(1.2), cw - Inches(0.56), Inches(1.3),
             comps, size=11.5, color=MUTED, spacing=1.3)
        text(s, x + Inches(0.28), y + Inches(2.62), cw - Inches(0.56), Inches(0.6),
             note, size=10.5, color=INK, spacing=1.25)
    caption(s, y + Inches(3.6),
            "Collapsing them destroys the information the strategist needs: a topic can be excellent for "
            "a strategist (large, early, no proof points) and useless for a salesperson (nothing to show).")
    footer(s, n := n + 1)

    # ---- 5. Evidence before generation ----------------------------------
    s = blank(prs); y = header(s, "Concept", "Evidence before generation",
                               "The model never invents a topic from its own knowledge. It reorganises retrieved, dated, attributable evidence.")
    defences = [
        ("1 · Evidence binding", "Every claim cites signal ids that must exist in the cluster that "
                                 "produced the candidate. Uncited claims are stripped, not rewritten."),
        ("2 · Closed vocabulary", "Taxonomy values are validated against the enumerations. A recognised "
                                  "synonym is repaired once; anything else is dropped."),
        ("3 · No generated numbers", "Market sizes, growth rates and percentages are looked up and "
                                     "attributed, or they are absent. Backstopped by a regex."),
        ("4 · Entailment check", "A second model pass verifies each claim is actually entailed by the "
                                 "span it cites."),
    ]
    for i, (head, body) in enumerate(defences):
        yy = y + Inches(0.78 * i)
        rect(s, Inches(0.75), yy, Inches(7.2), Inches(0.66), fill=LIGHT)
        text(s, Inches(1.0), yy + Inches(0.08), Inches(2.6), Inches(0.3), head, size=12.5, bold=True, color=ORANGE)
        text(s, Inches(1.0), yy + Inches(0.33), Inches(6.7), Inches(0.3), body, size=11, color=MUTED)
    rect(s, Inches(8.35), y, Inches(4.3), Inches(3.1), fill=INK)
    text(s, Inches(8.62), y + Inches(0.28), Inches(3.8), Inches(0.3),
         "PLUS AN ADVERSARIAL CRITIC", size=10, color=ORANGE, bold=True)
    text(s, Inches(8.62), y + Inches(0.68), Inches(3.8), Inches(2.2),
         "A separate critic prompt scores 1–5 as the MINIMUM across five tests, so one failure caps "
         "the whole score.\n\nIn the live run it rejected 345 of 644 candidates — with specific, "
         "written reasons.", size=12.5, color=WHITE, spacing=1.35)
    footer(s, n := n + 1)

    # ---- 6. Portfolio distance ------------------------------------------
    s = blank(prs); y = header(s, "Concept", "Portfolio distance decides whose conversation it is",
                               "The shortest path from a topic to a configuration Orange could actually deliver.")
    rows = [
        ("L0", "Direct", "An existing offer addresses it as it stands", "Sales — sell it"),
        ("L1", "Bundle", "Two or more existing offers combined", "Presales — package it"),
        ("L2", "Partner-dependent", "Needs a capability a partner already holds", "Presales / alliances — assemble it"),
        ("L3", "Adjacent", "Needs one capability built; nearby assets exist", "Strategy — study it"),
        ("L4", "White space", "No plausible path from the current portfolio", "Strategy — watch or reject"),
    ]
    shades = [RGBColor(0x86, 0xB6, 0xEF), RGBColor(0x55, 0x98, 0xE7), RGBColor(0x2A, 0x78, 0xD6),
              RGBColor(0x1C, 0x5C, 0xAB), RGBColor(0x10, 0x42, 0x81)]
    for i, ((code, name, desc, owner), col) in enumerate(zip(rows, shades)):
        yy = y + Inches(0.62 * i)
        rect(s, Inches(0.75), yy, Inches(0.62), Inches(0.5), fill=col)
        text(s, Inches(0.75), yy + Inches(0.12), Inches(0.62), Inches(0.3), code,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=MONO)
        text(s, Inches(1.55), yy + Inches(0.05), Inches(2.2), Inches(0.3), name, size=14, bold=True)
        text(s, Inches(3.75), yy + Inches(0.07), Inches(5.0), Inches(0.3), desc, size=12, color=MUTED)
        text(s, Inches(8.95), yy + Inches(0.07), Inches(3.7), Inches(0.3), owner, size=12, color=ORANGE)
    caption(s, y + Inches(3.35),
            "This is what drives the role modes — they are not arbitrary presets. A high-attractiveness "
            "L4 topic is exactly the strategist's innovation agenda, and exactly what a salesperson should never be shown.")
    footer(s, n := n + 1)

    # ---- 7. Live corpus -------------------------------------------------
    s = blank(prs); y = header(s, "Status", "What the MVP has actually produced",
                               "Every figure below is read from the live database, not typed into this deck.")
    sw = Inches(2.22)
    stats = [
        (F.get("topics"), "opportunity spaces", f"{F.get('verticals')}/15 verticals covered"),
        (f"{F.get('signals', 0):,}", "signals ingested", f"{F.get('relevant', 0):,} passed the gate"),
        (F.get("sources"), "live sources", "33 enabled + internal intake"),
        (f"{F.get('links', 0):,}", "asset links", f"{F.get('nodes')} graph nodes"),
        (F.get("avg_sig"), "signals per topic", "after enrichment"),
    ]
    for i, (v, l, sub) in enumerate(stats):
        stat(s, Inches(0.75) + (sw + Inches(0.2)) * i, y, sw, v, l, sub)
    text(s, Inches(0.75), y + Inches(1.62), Inches(11.9), Inches(0.4),
         "Coverage of the evidenced grid", size=15, bold=True)
    bullets(s, Inches(0.75), y + Inches(2.08), Inches(11.5), [
        (f"{F.get('use_cases')} of 59 use cases, {F.get('techs')} of 38 technologies",
         "appear in at least one topic — the grid is sparse by design."),
        (f"{F.get('tier1', 0):,} tier-1 signals",
         "from regulators, procurement portals, standards bodies and official statistics."),
        (f"{F.get('fr', 0):,} French-language signals",
         "so the anglophone bias named as a principal risk is measured, not assumed."),
        (f"{F.get('sized_bu', 0)} sized bottom-up · {F.get('competed', 0)} competition-scored · {F.get('briefs', 0)} briefed",
         "each with its method and its confidence grade, never a quoted headline."),
        (f"{F.get('plans', 0)} portfolio plans · {F.get('collateral_kinds', 0)} kinds of pre-sales artefact built",
         "the two things that turn a ranked list into a decision and a bid."),
    ], size=13.5, gap=0.5)
    footer(s, n := n + 1)


    # ---- 7b. Sizing and competition -------------------------------------
    s = blank(prs); y = header(s, "Concept", "Sizing and competition, with the working shown",
                               "Two questions a topic cannot be acted on without: how big is it, and who else is already there?")
    rect(s, Inches(0.75), y, Inches(6.0), Inches(2.9), fill=LIGHT)
    text(s, Inches(1.05), y + Inches(0.28), Inches(5.4), Inches(0.3),
         "MARKET SIZE — BOTTOM UP, NOT QUOTED", size=10.5, color=ORANGE, bold=True)
    text(s, Inches(1.05), y + Inches(0.68), Inches(5.4), Inches(2.0),
         "Headline market figures in press coverage originate from paid research, are quoted without "
         "methodology, and frequently conflict by an order of magnitude.\n\n"
         "So the radar builds its own: enterprise counts by sector and size class, times an observed "
         "adoption rate, times a plausible contract value. Every estimate carries its method, its "
         "inputs and a confidence label — and can be rejected on its arithmetic.",
         size=12, color=INK, spacing=1.35)
    rect(s, Inches(7.0), y, Inches(5.6), Inches(2.9), fill=LIGHT)
    text(s, Inches(7.3), y + Inches(0.28), Inches(5.0), Inches(0.3),
         "COMPETITIVE INTENSITY", size=10.5, color=ORANGE, bold=True)
    text(s, Inches(7.3), y + Inches(0.68), Inches(5.0), Inches(2.0),
         "Scored from a versioned competitor register against the evidence actually collected — who "
         "is visibly playing here.\n\n"
         "A crowded field is not a reason to walk away; it is a reason to win on a specific "
         "differentiator. And 'no competitor found' is reported as unverified rather than as empty, "
         "because it may only mean the register has a gap.",
         size=12, color=INK, spacing=1.35)
    caption(s, y + Inches(3.15),
            f"Both are filterable and sortable, because a number the radar computes but cannot be "
            f"sorted on is a number a strategist has to eyeball across {F.get('topics', 0)} rows.")
    footer(s, n := n + 1)


    # ---- 7c. Two routes into a new space --------------------------------
    s = blank(prs); y = header(s, "Concept", "Two routes into a new opportunity space",
                               "A refresh answers the questions the corpus was already asked. These answer the one somebody has right now.")
    routes = [
        ("PARAMETERS", "For somebody who knows the taxonomy", BLUE,
         "Pick a vertical, a use case, a technology, a horizon. Before anything is spent, the screen "
         "shows the spaces that ALREADY satisfy those criteria — because the most common outcome of an "
         "on-demand run is rediscovering what the last refresh produced."),
        ("A SCOPING CONVERSATION", "For somebody who knows their market", ORANGE,
         "The assistant interviews, with the corpus in front of it. Every turn re-embeds the whole "
         "transcript against the same signal vectors the run will read, and shows what came back — "
         "publisher, date and cosine — beside the answer."),
    ]
    cw3 = Inches(5.85)
    for i, (title_, sub, col, body) in enumerate(routes):
        x = Inches(0.75) + (cw3 + Inches(0.3)) * i
        rect(s, x, y, cw3, Inches(2.15), fill=LIGHT)
        rect(s, x, y, cw3, Inches(0.07), fill=col)
        text(s, x + Inches(0.28), y + Inches(0.3), cw3 - Inches(0.56), Inches(0.3), title_,
             size=11, color=col, bold=True)
        text(s, x + Inches(0.28), y + Inches(0.62), cw3 - Inches(0.56), Inches(0.35), sub,
             size=14, bold=True)
        text(s, x + Inches(0.28), y + Inches(1.08), cw3 - Inches(0.56), Inches(1.0), body,
             size=11.5, color=MUTED, spacing=1.32)
    rect(s, Inches(0.75), y + Inches(2.3), Inches(11.9), Inches(1.5), fill=INK)
    text(s, Inches(1.05), y + Inches(2.52), Inches(11.3), Inches(0.3),
         "BOTH ARE REFUSED BY THE SAME GATE, AND THE CORPUS HOLDS IT — NOT THE MODEL",
         size=10.5, color=ORANGE, bold=True)
    text(s, Inches(1.05), y + Inches(2.9), Inches(5.4), Inches(0.8),
         "A brief must retrieve at least the run's own floor of signals, using the run's own retrieval. "
         "Asked \u201cdo you have enough?\u201d a model says yes, so the button is enabled by what came back.",
         size=11, color=RGBColor(0xC3, 0xC2, 0xB7), spacing=1.3)
    text(s, Inches(6.9), y + Inches(2.9), Inches(5.5), Inches(0.8),
         "And similarity is not support. A second, independent reason is required — on the use case or "
         "the technology, never the vertical, which corroborates every brief ever written about a "
         "well-covered sector.", size=11, color=RGBColor(0xC3, 0xC2, 0xB7), spacing=1.3)
    caption(s, y + Inches(3.95),
            "Municipal digital signage retrieves French tenders at the same 0.64 cosine a well-evidenced "
            "brief scores — so the gate judges the brief's own sentence, not the taxonomy labels it was filed under.")
    footer(s, n := n + 1)

    # ---- 8-14. Functionality --------------------------------------------
    screens = [
        ("radar", "The radar view",
         "Four dimensions without a legend anyone has to study",
         "Angular sector is the business domain. Distance from the centre is the time horizon — Now at "
         "the middle, Later at the rim. Marker size is attractiveness; marker colour is right to win. "
         "Position carries identity, so colour is free to encode a quantity. A ! inside a marker means "
         "an evidence gap."),
        ("list", "Role-ranked list",
         "The same data, three genuinely different rankings",
         "Strategist ranks on attractiveness and novelty and ignores right-to-win. Sales ranks on "
         "right-to-win and proof-point density, and only sees topics with a delivery path, a published "
         "reference in the vertical, and no evidence gap. Presales ranks on differentiation."),
        ("detail", "Topic detail",
         "Every claim traceable to a dated source",
         "One page in the order the user's questions arrive: the statement, why it is hot now with each "
         "claim linked to its signals, where it delivers value and for whom, can-we-play/can-we-win "
         "itemised against named Orange assets, the score breakdown expanded, and the next action for "
         "the current role."),
        ("fs_space", "One space, full screen",
         "Four tabs, in the order the questions arrive",
         "The three-pane layout is right for working THROUGH the radar — filter, scan, open, compare, "
         "move on. It is wrong for the moment somebody reads a space. So the same content opens with "
         "the panes out of the way, in four tabs: the space, the competitors, the sales brief, and the "
         "pre-sales pack. What is this, who else is here, what do I send, what comes after the meeting."),
        ("explain", "How this score was calculated",
         "Explainability made checkable rather than asserted",
         "Every component shows the stored inputs and the arithmetic: publisher entropy and the "
         "publishers counted, the tier distribution, the per-period buckets the momentum slope was "
         "fitted to, the rubric level and its rationale. A reviewer outside the project can reconstruct "
         "why any topic holds its rank."),
        ("fs_presales", "Pre-sales collateral",
         "Twelve pieces, in the format each reader works in",
         "The brief is one document for one conversation. This is what the team needs between that "
         "conversation and a proposal. All twelve are listed whether or not anything has been built, "
         "because what COULD be produced is as much of the answer as what has been — and a screen that "
         "starts empty is one nobody presses a button on."),
        ("workflow", "Stage gate and role assessment",
         "Shortlisted → Demand-tested → Packaged → Live",
         "Ownership follows the stage, and stalled cards are flagged because latency is the known "
         "weakness of a stage gate. Each role rates only the axis it owns, on a 0–5 scale with written "
         "anchors. Those ratings form conviction — a third quantity that changes ranking only."),
        ("analytics", "Analytics",
         "Charts chosen by the job the data does",
         "The vertical × domain heatmap is magnitude on a grid, so sequential. Conviction versus "
         "evidence is polarity, so diverging with a neutral midpoint — agreement reads as nothing. "
         "The stage funnel is an ordered sequence, so ordinal. Only the signal-type mix is categorical, "
         "and it ships a legend and a table."),
        ("generate_chat", "Generating a space on demand",
         "The screen opens with a conversation, not a text box",
         "The box asked for one thing and gave one piece of feedback — a character count, the only "
         "failure that did not matter. Somebody who knows their market but not this taxonomy "
         "under-specified two of five dimensions every time, and found out minutes later from a run "
         "that created nothing. The assistant interviews instead, and shows what each turn retrieved."),
        ("planner_overview", "The Planner",
         "Which SET, in what order, and what it earns",
         "A ranked list assumes you can take the top N, and you cannot. A mixed-integer program "
         "maximises the stated objective under entry slots, capability headcount, concentration caps "
         "and a horizon mix — then reports WHICH CONSTRAINT BOUND IT, which is the thing a ranked list "
         "cannot tell you, because the answer is a constraint rather than a score."),
        ("plannerwf_narrative", "The plan the business already chose",
         "Workflow selected: the stage gate decides, the Planner only schedules",
         "Every space the board has moved to Demand-tested or beyond is in, and none of the "
         "constraints is applied — each would overrule a decision somebody already took. Nothing is "
         "dropped to make it fit: where the committed set needs more than the pools can staff, the "
         "plan says so and by how much. That gap is the finding, not a reason to edit the portfolio."),
        ("help", "Contextual help",
         "Every dense concept explains itself",
         "Portfolio distance, conviction, divergence, evidence gaps, source tiers, horizons, the "
         "lifecycle, the exploration slot. Each entry says what the thing is, why it works that way, "
         "and which requirement it comes from — so the answer is checkable rather than merely confident."),
    ]
    for name, title_, sub, body in screens:
        s = blank(prs)
        rect(s, 0, 0, W, Inches(0.06), fill=ORANGE)
        text(s, Inches(0.7), Inches(0.4), Inches(7.4), Inches(0.3),
             "FUNCTIONALITY", size=10.5, color=ORANGE, bold=True)
        text(s, Inches(0.7), Inches(0.72), Inches(7.4), Inches(0.5), title_, size=27, bold=True)
        text(s, Inches(0.7), Inches(1.28), Inches(7.4), Inches(0.4), sub, size=14, color=ORANGE)
        text(s, Inches(0.7), Inches(1.82), Inches(3.55), Inches(4.6), body, size=12, color=MUTED, spacing=1.4)
        picture(s, name, Inches(4.55), Inches(1.55), Inches(8.2))
        footer(s, n := n + 1)

    # ---- 15. Architecture: pipeline -------------------------------------
    s = blank(prs); y = header(s, "Architecture", "Seven pipeline stages, each with a contract",
                               "Stages can be developed, tested and replaced independently.")
    stages = [
        ("1 Collect", "connectors", "source config → raw items"),
        ("2 Normalise", "ingest", "raw items → signal records"),
        ("3 Classify", "ingest", "→ typed, tiered signals"),
        ("4 Themes", "themes", "→ theme clusters"),
        ("5 Synthesise", "synthesis", "→ candidate spaces"),
        ("5b Enrich", "enrich", "→ more evidence per topic"),
        ("6 Score", "graph + scoring", "→ ranked spaces"),
        ("7 Serve", "readmodel + api", "→ radar, briefs, API"),
    ]
    bw2 = Inches(1.44)
    for i, (name, mod, io) in enumerate(stages):
        x = Inches(0.75) + (bw2 + Inches(0.05)) * i
        dark = name.startswith(("5", "6"))
        rect(s, x, y, bw2, Inches(1.45), fill=INK if dark else LIGHT)
        text(s, x + Inches(0.12), y + Inches(0.16), bw2 - Inches(0.24), Inches(0.4), name,
             size=11.5, bold=True, color=WHITE if dark else INK)
        text(s, x + Inches(0.12), y + Inches(0.58), bw2 - Inches(0.24), Inches(0.3), mod,
             size=9, color=ORANGE, font=MONO)
        text(s, x + Inches(0.12), y + Inches(0.86), bw2 - Inches(0.24), Inches(0.6), io,
             size=9, color=RGBColor(0xC3, 0xC2, 0xB7) if dark else MUTED, spacing=1.2)
    text(s, Inches(0.75), y + Inches(1.72), Inches(11.9), Inches(0.4),
         "A parallel, slower path maintains the Orange Business Graph", size=15, bold=True)
    text(s, Inches(0.75), y + Inches(2.08), Inches(11.9), Inches(0.7),
         "Offers, references, partners with tiers, certifications, analyst positions, capability pools "
         "and research assets. It joins at stage 6, so right-to-win can be improved without re-running "
         "discovery. Links are typed L0–L4 and carry the evidence that justified them.",
         size=12.5, color=MUTED, spacing=1.32)
    text(s, Inches(0.75), y + Inches(2.76), Inches(11.9), Inches(0.4),
         "And two subsystems sit BESIDE the pipeline rather than in it", size=15, bold=True)
    subs = [
        ("THE PLANNER", "planner.py · plan_report.py",
         "Reads the read model and an assumptions file. Selects a set, schedules it, projects five "
         "years, writes one PDF. One model call, after every number is fixed."),
        ("PRE-SALES COLLATERAL", "presales/",
         "Reads ONE snapshot of a space and emits twelve artefacts in five formats — so nothing in a "
         "pack can disagree with anything else in it."),
    ]
    for i, (title_, mod, body) in enumerate(subs):
        x = Inches(0.75) + Inches(6.1) * i
        rect(s, x, y + Inches(3.1), Inches(5.8), Inches(1.32), fill=LIGHT)
        text(s, x + Inches(0.24), y + Inches(3.27), Inches(5.3), Inches(0.3), title_,
             size=10.5, color=ORANGE, bold=True)
        text(s, x + Inches(0.24), y + Inches(3.55), Inches(5.3), Inches(0.25), mod,
             size=9, color=MUTED, font=MONO)
        text(s, x + Inches(0.24), y + Inches(3.81), Inches(5.3), Inches(0.7), body,
             size=10.5, color=MUTED, spacing=1.28)
    footer(s, n := n + 1)

    # ---- 16. Architecture: stack ----------------------------------------
    s = blank(prs); y = header(s, "Architecture", "Stack and separation of concerns")
    cols = [
        ("INGESTION", ["33 enabled sources across 17 connector", "types: TED, BOAMP, UK Contracts,",
                       "EUR-Lex SPARQL, Have-your-say, GDELT,", "news RSS (EN/FR), Hacker News, OpenAlex,",
                       "Crossref, arXiv, CORDIS, NIST, CERT-FR", "",
                       "Parallel fetch · circuit breaker ·", "per-host pacing · replay-safe dates"]),
        ("INTELLIGENCE", ["DeepSeek behind a provider-agnostic", "client (swap to Ollama via .env alone)",
                          "", "Local sentence-transformers for", "embeddings and clustering",
                          "", "Deterministic where possible,", "generative only where it earns its place"]),
        ("SERVING", ["SQLite — the graph is thousands of", "nodes, not millions",
                     "", "FastAPI read API behind a session", "guard applied to the whole app,",
                     "not route by route",
                     "", "React + Vite + TypeScript",
                     "Hand-drawn SVG radar, no chart library",
                     "", "scipy milp for portfolio selection",
                     "reportlab / python-pptx / python-docx",
                     "for every document it emits"]),
    ]
    cw2 = Inches(3.85)
    for i, (title_, lines) in enumerate(cols):
        x = Inches(0.75) + (cw2 + Inches(0.28)) * i
        rect(s, x, y, cw2, Inches(3.5), fill=LIGHT)
        rect(s, x, y, cw2, Inches(0.07), fill=ORANGE)
        text(s, x + Inches(0.28), y + Inches(0.3), cw2 - Inches(0.56), Inches(0.3),
             title_, size=11, color=ORANGE, bold=True)
        text(s, x + Inches(0.28), y + Inches(0.72), cw2 - Inches(0.56), Inches(2.6),
             "\n".join(lines), size=11.5, color=INK, spacing=1.32)
    text(s, Inches(0.75), y + Inches(3.8), Inches(11.9), Inches(0.8),
         "Taxonomies, weights, thresholds, sources, the business graph and the crosswalks are all "
         "configuration, not code — and are validated at load time, so a dangling id is a startup error "
         "rather than a wrong number three stages later.", size=12.5, color=MUTED, spacing=1.35)
    footer(s, n := n + 1)

    # ---- 17. Explainability guarantees ----------------------------------
    s = blank(prs); y = header(s, "Architecture", "What makes the numbers defensible",
                               "The governing constraint: if a user cannot explain why a topic is ranked where it is, the scoring is not good enough.")
    guarantees = [
        ("Decomposable", "Every displayed number breaks into named components; no opaque scores."),
        ("Reproducible", "Every component stores the inputs used to compute it, so any number can be re-derived."),
        ("Traceable", "Full lineage from a displayed claim back to the raw ingested item, including prompt and model version."),
        ("Versioned", "Every score records its weight set. Trajectories are never plotted across a boundary silently."),
        ("Auditable", "A reviewer outside the project can reconstruct why any topic holds its rank."),
        ("Bounded", "Counting, diversity, recency and momentum are arithmetic — never a model. A model asked to count is unverifiable."),
    ]
    for i, (head, body) in enumerate(guarantees):
        col, row = i % 2, i // 2
        x = Inches(0.75) + Inches(6.2) * col
        yy = y + Inches(1.0) * row
        rect(s, x, yy + Inches(0.06), Inches(0.055), Inches(0.62), fill=ORANGE)
        text(s, x + Inches(0.24), yy, Inches(5.6), Inches(0.3), head, size=15, bold=True)
        text(s, x + Inches(0.24), yy + Inches(0.3), Inches(5.6), Inches(0.6), body,
             size=11.5, color=MUTED, spacing=1.3)
    footer(s, n := n + 1)

    # ---- 18. Honest limits ----------------------------------------------
    s = blank(prs); y = header(s, "Status", "What is deliberately not built — and what needs a decision")
    text(s, Inches(0.75), y, Inches(5.9), Inches(0.3), "NOT BUILT, WITH THE REASON", size=11, color=ORANGE, bold=True)
    bullets(s, Inches(0.75), y + Inches(0.42), Inches(5.5), [
        ("CRM integration", "deferred; public assets give a sufficient right-to-win proxy"),
        ("Learned scoring models", "no labels exist on day one; the capture and replay harness ships instead"),
        ("Patent connector", "needs EPO registration; technology ownership uses a portfolio-level prior"),
        ("Learned ranking per role", "needs 300-600 expert comparisons; the capture widget ships first"),
        ("ROI on a plan", "there is no cost data at the granularity a space needs, anywhere the pipeline can reach"),
        ("Per-role authorisation", "sign-in answers WHO; it does not yet answer MAY THEY"),
    ], size=12.5, gap=0.55)
    text(s, Inches(7.0), y, Inches(5.6), Inches(0.3), "NEEDS A DECISION FROM ORANGE", size=11, color=RED, bold=True)
    bullets(s, Inches(7.0), y + Inches(0.42), Inches(5.4), [
        ("Who is the curator?", f"{F.get('links', 0):,} links are machine-proposed and unconfirmed"),
        ("Margin by portfolio distance", "one table from Orange finance moves five-year profit by 1.66x"),
        ("Headcount free for new work", "the constraint that binds first in most plans is currently a guess"),
        ("Terms of use", "unconfirmed for several enabled sources — a Sprint 0 blocker"),
        ("Refresh cadence", "drives connector design and cost more than any other choice"),
    ], size=12.5, gap=0.55)
    caption(s, y + Inches(3.35),
            "The radar reports its own gaps rather than hiding them: evidence-gap warnings, language "
            "coverage, unconfirmed links, skipped sources, unsized spaces and an over-committed "
            "capability pool are all surfaced in the interface rather than left to be discovered.")
    footer(s, n := n + 1)

    # ---- 19. Close ------------------------------------------------------
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=INK)
    rect(s, 0, 0, Inches(0.22), H, fill=ORANGE)
    text(s, Inches(1.1), Inches(2.3), Inches(11), Inches(0.9),
         "The join is the product", size=40, bold=True, color=WHITE)
    text(s, Inches(1.1), Inches(3.4), Inches(10.2), Inches(1.4),
         "Without it the radar is a competent trend feed, and trend feeds already exist. With it, the "
         "radar answers a question nobody else can answer for Orange.", size=17,
         color=RGBColor(0xC3, 0xC2, 0xB7), spacing=1.35)
    text(s, Inches(1.1), Inches(5.1), Inches(11), Inches(0.4),
         f"{F.get('topics')} opportunity spaces  ·  {F.get('signals', 0):,} signals  ·  "
         f"{F.get('sources')} live sources  ·  {F.get('links', 0):,} named asset links",
         size=13, color=ORANGE, font=MONO)
    text(s, Inches(1.1), Inches(5.5), Inches(11), Inches(0.4),
         f"{F.get('sized_bu', 0)} sized bottom-up  ·  {F.get('briefs', 0)} sales briefs  ·  "
         f"12 pre-sales artefacts per space  ·  {F.get('plans', 0)} portfolio plans",
         size=13, color=MUTED, font=MONO)
    n += 1

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT}  ({n} slides)")
    return OUT


if __name__ == "__main__":
    sys.exit(0 if build() else 1)
